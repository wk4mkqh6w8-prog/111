import os
import json
import hmac
import logging
import asyncio
import hashlib
import subprocess
import threading
import time
import re
import html
import textwrap
from types import SimpleNamespace
from datetime import datetime, timedelta

import base64
import shutil
import tempfile
from pathlib import Path
from pypdf import PdfReader  # pip install PyPDF2
from gtts import gTTS  # ← добавь импорт рядом с остальными
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import httpx
import requests
from dotenv import load_dotenv
from fastapi import FastAPI, Request
from fastapi.responses import HTMLResponse, PlainTextResponse
import uvicorn
from fpdf import FPDF

from openai import OpenAI
from telegram import Update, InlineKeyboardMarkup, InlineKeyboardButton, ReplyKeyboardMarkup
from telegram.ext import (
    Application, ApplicationBuilder,
    CommandHandler, CallbackQueryHandler, MessageHandler,
    ContextTypes, filters,
)

# =========================
# Конфиг и клиенты
# =========================
load_dotenv()

BOT_TOKEN      = os.getenv("BOT_TOKEN", "")
OPENAI_KEY     = os.getenv("OPENAI_KEY", "")
# Поддержка пула ключей: можно задать OPENAI_KEYS="sk-1,sk-2,sk-3"
OPENAI_KEYS_RAW = os.getenv("OPENAI_KEYS", "")
OPENAI_KEYS = [k.strip() for k in (OPENAI_KEYS_RAW or OPENAI_KEY or "").split(",") if k and k.strip()]
DEEPSEEK_KEY   = os.getenv("DEEPSEEK_KEY", "")
CRYPTOPAY_KEY  = os.getenv("CRYPTOPAY_KEY", "")
REPLICATE_KEY  = os.getenv("REPLICATE_KEY", "")
ADMIN_ID       = int(os.getenv("ADMIN_ID", "0"))
PORT           = int(os.getenv("PORT", "10000"))
SUPPORT_EMAIL      = os.getenv("SUPPORT_EMAIL", "support@neurobotgpt.ru")
PUBLIC_OFFER_URL   = os.getenv("PUBLIC_OFFER_URL", "https://disk.yandex.ru/i/wdHQVfYcJGjwhw")
SUPPORT_WORK_HOURS = os.getenv("SUPPORT_WORK_HOURS", "10:00–19:00 MSK")
PHOTO_COOLDOWN_SEC = int(os.getenv("PHOTO_COOLDOWN_SEC", "60"))  # КД на фото для всех, сек

if not BOT_TOKEN:
    raise RuntimeError("BOT_TOKEN пуст")
if not OPENAI_KEYS:
    raise RuntimeError("OPENAI_KEYS/OPENAI_KEY пуст")

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s %(levelname)s %(name)s: %(message)s"
)
logger = logging.getLogger("neurobot")

# Модели (реальные — для движка)
MODEL_OPENAI   = "OpenAI · GPT-4o-mini"
MODEL_DEEPSEEK = "DeepSeek · Chat"
DEFAULT_MODEL  = MODEL_OPENAI

# Выбор пользователя
_user_model_visual: dict[int, str] = {}  # «название модели» которое видит пользователь
_user_model: dict[int, str] = {}         # фактический backend (OpenAI/DeepSeek)
_awaiting_img_prompt: dict[int, bool] = {}
_pending_chat_rename: dict[int, int] = {}  # user_id -> chat_id
_last_answer: dict[int, str] = {}           # последний текстовый ответ для TTS
_long_reply_queue: dict[int, list[str]] = {}  # очереди «показать ещё»
_photo_cd_until: dict[int, float] = {}  # user_id -> unix timestamp до которого фото нельзя слать
_user_profiles: dict[int, dict[str, str]] = {}
_last_user_prompt: dict[int, str] = {}

PROFILE_STYLES = {
    "standard": "Стандарт",
    "friendly": "Дружелюбный",
    "formal": "Официальный",
    "expert": "Экспертный",
}

PROFILE_STYLE_INSTRUCTIONS = {
    "standard": "",
    "friendly": "Adopt a warm, encouraging tone and add light emoji where it improves clarity.",
    "formal": "Use a formal, professional tone with complete sentences.",
    "expert": "Respond like a subject-matter expert, referencing best practices and terminology.",
}

PROFILE_LANGUAGES = {
    "auto": "Авто",
    "ru": "Русский",
    "en": "English",
}

PROFILE_LANGUAGE_INSTRUCTIONS = {
    "auto": "Match the user's language. If unsure, default to Russian.",
    "ru": "Respond in Russian.",
    "en": "Respond in English.",
}

PROFILE_FORMATS = {
    "plain": "Обычный текст",
    "bullets": "Списки",
    "markdown": "Markdown",
}

PROFILE_FORMAT_INSTRUCTIONS = {
    "plain": "",
    "bullets": "Format the answer as concise bullet points.",
    "markdown": "Use clear Markdown formatting with headings and lists where helpful.",
}

PROFILE_THEMES = {
    "auto": "Авто",
    "light": "Светлая карточка",
    "dark": "Тёмная карточка",
}

PROFILE_THEME_INSTRUCTIONS = {
    "auto": "",
    "light": "Keep the tone upbeat and add a short positive closing.",
    "dark": "Use a slightly more atmospheric tone suitable for dark UI cards.",
}

QUICK_COMMANDS_KEYBOARD = ReplyKeyboardMarkup(
    [["/help", "/img"], ["/ppt", "/favorites"], ["/settings"]],
    resize_keyboard=True,
    selective=True,
)


async def _ensure_profile(user_id: int) -> dict[str, str]:
    profile = _user_profiles.get(user_id)
    if profile is None:
        profile = await get_user_profile_settings(user_id)
        _user_profiles[user_id] = profile
    return profile


def _profile_snapshot(user_id: int) -> dict[str, str]:
    profile = _user_profiles.get(user_id)
    if not profile:
        return dict(DEFAULT_PROFILE)
    return {
        "style": profile.get("style", DEFAULT_PROFILE["style"]),
        "language": profile.get("language", DEFAULT_PROFILE["language"]),
        "output_format": profile.get("output_format", DEFAULT_PROFILE["output_format"]),
        "theme": profile.get("theme", DEFAULT_PROFILE["theme"]),
    }


def _update_profile_cache(user_id: int, field: str, value: str):
    profile = _user_profiles.setdefault(user_id, dict(DEFAULT_PROFILE))
    profile[field] = value

# РЕЖИМЫ (ярлыки): реально влияют на подсказку
TASK_MODES = {
    "default": {
        "label": "Стандарт",
        "system": (
            "You are a helpful, concise assistant. Prefer clear steps and short answers unless "
            "the user asks for depth."
        ),
    },
    "coding": {
        "label": "Кодинг",
        "system": (
            "You are a senior software engineer. Provide runnable code with comments, point out pitfalls, "
            "and show minimal examples. Prefer Python/JS unless the user specifies otherwise."
        ),
    },
    "seo": {
        "label": "SEO",
        "system": (
            "You are an SEO strategist. Produce keyword-rich but natural copy, suggest title/H1/meta, "
            "and include semantic clusters and internal linking ideas when useful."
        ),
    },
    "translate": {
        "label": "Перевод",
        "system": (
            "You are a professional translator (RU↔EN). Preserve meaning, tone, and idioms. "
            "If the source is ambiguous, offer the two best variants."
        ),
    },
    "summarize": {
        "label": "Резюме",
        "system": (
            "You are a world-class summarizer. Output structured bullet points, key facts, and action items. "
            "Keep it brief unless asked to expand."
        ),
    },
    "creative": {
        "label": "Креатив",
        "system": (
            "You are a creative copywriter. Offer punchy hooks, strong voice, and multiple variants when helpful. "
            "Avoid clichés."
        ),
    },
}
_user_task_mode: dict[int, str] = {}  # хранит ключ режима пользователя

# ----- OpenAI clients pool + failover logic -----
from collections import deque

_oai_clients: dict[str, OpenAI] = {}
_openai_keys_ring = deque(OPENAI_KEYS)
_key_cooldowns: dict[str, float] = {}   # key -> unix_timestamp до какого молчим

from collections import deque

_recent_updates = deque(maxlen=1000)
_recent_set = set()

def _get_client(api_key: str) -> OpenAI:
    cli = _oai_clients.get(api_key)
    if cli is None:
        cli = OpenAI(api_key=api_key)
        _oai_clients[api_key] = cli
    return cli

def _mark_cooldown(api_key: str, seconds: int):
    _key_cooldowns[api_key] = time.time() + max(1, seconds)

def _pick_next_key() -> str | None:
    """Берём следующий ключ, который не на кулдауне."""
    if not _openai_keys_ring:
        return None
    now = time.time()
    for _ in range(len(_openai_keys_ring)):
        k = _openai_keys_ring[0]
        _openai_keys_ring.rotate(-1)
        if _key_cooldowns.get(k, 0) <= now:
            return k
    return None  # все в кулдауне

def _oai_chat_call(messages: list[dict], model: str, temperature: float = 0.7) -> str:
    """
    Вызывает chat.completions с автоматическим переключением ключей.
    Возвращает текст ответа или кидает RuntimeError, если все ключи не сработали.
    """
    last_err: Exception | None = None
    tried: set[str] = set()

    for _ in range(len(OPENAI_KEYS)):
        api_key = _pick_next_key()
        if not api_key or api_key in tried:
            break
        tried.add(api_key)
        client = _get_client(api_key)
        try:
            r = client.chat.completions.create(
                model=model,
                messages=messages,
                temperature=temperature,
            )
            return r.choices[0].message.content
        except Exception as e:
            # классифицируем и ставим разумный кулдаун
            status = getattr(e, "status_code", None)
            if status == 401:          # невалидный/отключённый ключ
                _mark_cooldown(api_key, 600)
            elif status in (429, 500, 503):
                _mark_cooldown(api_key, 60)    # лимит/перегруз/апстрим
            else:
                _mark_cooldown(api_key, 10)
            last_err = e
            continue

    raise RuntimeError(f"All OpenAI keys failed: {last_err!s}")

# =========================
# DB helpers
# =========================
from db import (  # noqa
    init_db, add_user, is_premium, can_send_message, set_premium,
    get_usage_today, get_free_credits, consume_free_credit, add_free_credits,
    set_referrer_if_empty, count_paid_users_today, count_paid_users_total,
    get_premium_expires, list_expired_unnotified, mark_expired_notified,
    revoke_premium,
    # новые:
    get_chat_mode, set_chat_mode, create_chat, list_chats,
    set_active_chat, get_active_chat, add_chat_message, get_chat_history,
    rename_chat, delete_chat,
    get_user_profile_settings, set_user_profile_value, DEFAULT_PROFILE,
    add_favorite_prompt, list_favorite_prompts, get_favorite_prompt, delete_favorite_prompt,
    set_chat_pinned, create_chat_share, get_chat_share, cleanup_chat_shares,
    get_chat_history_all
)

# =========================
# FastAPI & PTB
# =========================
app = FastAPI(title="NeuroBot API")
application: Application | None = None
_public_url: str | None = None
_keepalive_stop = threading.Event()

REF_BONUS   = 25
DAILY_LIMIT = 5
# --- Цены ---
PRICE_RUB = 500                 # цена для пользователя (в рублях)
PRICE_USDT = "5"                # сумма счёта для Crypto Pay (USDT), строкой как требует API
PRICE_RUB_TEXT = f"{PRICE_RUB} ₽"

# --- Диалоговые режимы ---
DIALOG_SIMPLE = "simple"  # Быстрые ответы (без памяти)
DIALOG_ROOMS  = "rooms"   # Диалоги с контекстом (чаты)

# ---------- LLM ----------
def _compose_prompt(user_id: int, user_text: str, profile: dict[str, str] | None = None) -> list[dict]:
    """Собираем сообщения с учётом выбранного режима и профиля пользователя."""
    mode_key = _user_task_mode.get(user_id, "default")
    sys_text = TASK_MODES.get(mode_key, TASK_MODES["default"])["system"]

    profile = profile or _profile_snapshot(user_id)
    instructions: list[str] = []

    style = profile.get("style", "standard")
    language = profile.get("language", "auto")
    output_format = profile.get("output_format", "plain")
    theme = profile.get("theme", "auto")

    if PROFILE_STYLE_INSTRUCTIONS.get(style):
        instructions.append(PROFILE_STYLE_INSTRUCTIONS[style])
    if PROFILE_LANGUAGE_INSTRUCTIONS.get(language):
        instructions.append(PROFILE_LANGUAGE_INSTRUCTIONS[language])
    if PROFILE_FORMAT_INSTRUCTIONS.get(output_format):
        instructions.append(PROFILE_FORMAT_INSTRUCTIONS[output_format])
    if PROFILE_THEME_INSTRUCTIONS.get(theme):
        instructions.append(PROFILE_THEME_INSTRUCTIONS[theme])

    if instructions:
        sys_text = f"{sys_text} {' '.join(instructions)}"

    return [
        {"role": "system", "content": sys_text},
        {"role": "user", "content": user_text},
    ]

def _ask_openai(user_id: int, prompt: str) -> str:
    profile = _profile_snapshot(user_id)
    msgs = _compose_prompt(user_id, prompt, profile)
    # используем failover wrapper
    return _oai_chat_call(messages=msgs, model="gpt-4o-mini", temperature=0.7)

def _ask_deepseek(user_id: int, prompt: str) -> str:
    if not DEEPSEEK_KEY:
        return "DeepSeek недоступен: не задан DEEPSEEK_KEY."
    try:
        import httpx
        url = "https://api.deepseek.com/chat/completions"
        headers = {"Authorization": f"Bearer {DEEPSEEK_KEY}", "Content-Type": "application/json"}
        payload = {
            "model": "deepseek-chat",
            "messages": _compose_prompt(user_id, prompt, _profile_snapshot(user_id)),
            "temperature": 0.7,
        }
        with httpx.Client(timeout=30) as s:
            resp = s.post(url, headers=headers, json=payload)
            if resp.status_code != 200:
                try:
                    err = resp.json()
                    msg = err.get("error", {}).get("message") or err.get("message") or str(err)
                except Exception:
                    msg = resp.text[:400]
                return f"DeepSeek API error {resp.status_code}: {msg}"
            data = resp.json()
        choice = (data or {}).get("choices", [{}])[0]
        msg = (choice or {}).get("message", {})
        text = msg.get("content") or (choice or {}).get("text") or ""
        return text or "DeepSeek вернул пустой ответ."
    except Exception as e:
        return f"Ошибка DeepSeek: {e!s}"

def ask_llm(user_id: int, prompt: str) -> str:
    real = _user_model.get(user_id, DEFAULT_MODEL)
    if real == MODEL_DEEPSEEK:
        return _ask_deepseek(user_id, prompt)
    return _ask_openai(user_id, prompt)

def ask_llm_context(user_id: int, history: list[tuple[str, str]], user_text: str) -> str:
    """
    history: список (role, content), роли: 'system' | 'user' | 'assistant'
    """
    # системное сообщение — как в обычном режиме (учитываем TASK_MODES):
    sys_text = TASK_MODES.get(_user_task_mode.get(user_id, "default"), TASK_MODES["default"])["system"]
    msgs = [{"role": "system", "content": sys_text}]
    for role, content in history:
        if role in ("user", "assistant"):
            msgs.append({"role": role, "content": content})
    msgs.append({"role": "user", "content": user_text})

    profile = _profile_snapshot(user_id)
    instructions: list[str] = []
    style = profile.get("style")
    language = profile.get("language")
    output_format = profile.get("output_format")
    theme = profile.get("theme")
    if PROFILE_STYLE_INSTRUCTIONS.get(style):
        instructions.append(PROFILE_STYLE_INSTRUCTIONS[style])
    if PROFILE_LANGUAGE_INSTRUCTIONS.get(language):
        instructions.append(PROFILE_LANGUAGE_INSTRUCTIONS[language])
    if PROFILE_FORMAT_INSTRUCTIONS.get(output_format):
        instructions.append(PROFILE_FORMAT_INSTRUCTIONS[output_format])
    if PROFILE_THEME_INSTRUCTIONS.get(theme):
        instructions.append(PROFILE_THEME_INSTRUCTIONS[theme])
    if instructions:
        msgs[0]["content"] = f"{msgs[0]['content']} {' '.join(instructions)}"

    real = _user_model.get(user_id, DEFAULT_MODEL)
    if real == MODEL_DEEPSEEK:
        # DeepSeek
        try:
            import httpx
            url = "https://api.deepseek.com/chat/completions"
            headers = {"Authorization": f"Bearer {DEEPSEEK_KEY}", "Content-Type": "application/json"}
            payload = {"model": "deepseek-chat", "messages": msgs, "temperature": 0.7}
            with httpx.Client(timeout=30) as s:
                resp = s.post(url, headers=headers, json=payload)
                if resp.status_code != 200:
                    try:
                        err = resp.json()
                        msg = err.get("error", {}).get("message") or err.get("message") or str(err)
                    except Exception:
                        msg = resp.text[:400]
                    return f"DeepSeek API error {resp.status_code}: {msg}"
                data = resp.json()
            choice = (data or {}).get("choices", [{}])[0]
            m = (choice or {}).get("message", {})
            text = m.get("content") or (choice or {}).get("text") or ""
            return text or "DeepSeek вернул пустой ответ."
        except Exception as e:
            return f"Ошибка DeepSeek: {e!s}"
    else:
        # OpenAI
        # Используем общий вызов с переключением ключей
        profile = _profile_snapshot(user_id)
        instructions = []
        style = profile.get("style")
        language = profile.get("language")
        output_format = profile.get("output_format")
        theme = profile.get("theme")
        if PROFILE_STYLE_INSTRUCTIONS.get(style, ""):
            instructions.append(PROFILE_STYLE_INSTRUCTIONS[style])
        if PROFILE_LANGUAGE_INSTRUCTIONS.get(language, ""):
            instructions.append(PROFILE_LANGUAGE_INSTRUCTIONS[language])
        if PROFILE_FORMAT_INSTRUCTIONS.get(output_format, ""):
            instructions.append(PROFILE_FORMAT_INSTRUCTIONS[output_format])
        if PROFILE_THEME_INSTRUCTIONS.get(theme, ""):
            instructions.append(PROFILE_THEME_INSTRUCTIONS[theme])

        if instructions:
            msgs[0]["content"] = f"{msgs[0]['content']} {' '.join(instructions)}"
        return _oai_chat_call(messages=msgs, model="gpt-4o-mini", temperature=0.7)

def _transcribe_audio_file_sync(path: Path) -> str:
    """
    Выполняет синхронную расшифровку аудио через OpenAI (используем пул ключей).
    Возвращает распознанный текст.
    """
    last_err: Exception | None = None
    tried: set[str] = set()

    for _ in range(len(OPENAI_KEYS)):
        api_key = _pick_next_key()
        if not api_key or api_key in tried:
            break
        tried.add(api_key)
        client = _get_client(api_key)
        try:
            with path.open("rb") as audio_file:
                result = client.audio.transcriptions.create(
                    model="gpt-4o-transcribe",
                    file=audio_file,
                    response_format="text",
                )
            if isinstance(result, str):
                return result.strip()
            text = getattr(result, "text", "")
            if text:
                return str(text).strip()
            return ""
        except Exception as e:
            status = getattr(e, "status_code", None)
            if status == 401:
                _mark_cooldown(api_key, 600)
            elif status in (429, 500, 503):
                _mark_cooldown(api_key, 60)
            else:
                _mark_cooldown(api_key, 10)
            last_err = e
            continue

    raise RuntimeError(f"Transcription failed: {last_err!s}")

async def tts_and_send(user_id: int, chat_id: int, text: str, bot):
    """Озвучивает text через gTTS и пытается отправить голосовое (OPUS)."""
    tmpdir = Path(tempfile.gettempdir())
    mp3_path = tmpdir / f"tts_{user_id}_{int(time.time())}.mp3"
    ogg_path = mp3_path.with_suffix(".ogg")

    try:
        # gTTS ограничений по квоте нет; режем текст на всякий случай
        tts = gTTS(text=text[:4000], lang="ru")
        tts.save(str(mp3_path))

        sent_voice = False
        try:
            if await _convert_mp3_to_ogg(mp3_path, ogg_path):
                with open(ogg_path, "rb") as voice_file:
                    await bot.send_voice(
                        chat_id=chat_id,
                        voice=voice_file,
                        caption="Озвучено 🎧",
                    )
                sent_voice = True
        except Exception as convert_err:
            logger.warning("TTS voice conversion failed: %s", convert_err)

        if not sent_voice:
            # Fallback: отправляем MP3 как обычное аудио.
            with open(mp3_path, "rb") as audio_file:
                await bot.send_audio(
                    chat_id=chat_id,
                    audio=audio_file,
                    caption="Озвучено 🎧",
                    title="TTS",
                    filename=mp3_path.name,
                )
    except Exception as e:
        await bot.send_message(chat_id=chat_id, text=f"Не вышло озвучить: {e}")
    finally:
        for path in (mp3_path, ogg_path):
            try:
                path.unlink(missing_ok=True)
            except Exception:
                pass


async def _convert_mp3_to_ogg(mp3_path: Path, ogg_path: Path) -> bool:
    """
    Конвертирует MP3 в OGG/OPUS через ffmpeg (если доступен).
    Возвращает True при успехе.
    """
    ffmpeg_bin = shutil.which("ffmpeg")
    if not ffmpeg_bin:
        return False

    cmd = [
        ffmpeg_bin,
        "-y",
        "-loglevel",
        "error",
        "-i",
        str(mp3_path),
        "-ac",
        "1",
        "-ar",
        "48000",
        "-c:a",
        "libopus",
        "-b:a",
        "48k",
        "-vbr",
        "on",
        str(ogg_path),
    ]

    def _run_ffmpeg():
        return subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE)

    proc = await asyncio.to_thread(_run_ffmpeg)
    return proc.returncode == 0 and ogg_path.exists()

async def on_tts_btn(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass
    uid = q.from_user.id
    text = _last_answer.get(uid)
    if not text:
        await q.message.reply_text("Нет текста для озвучки.")
        return
    await tts_and_send(uid, q.message.chat_id, text, context.bot)

# =========================
# Хелперы для фото/доков
# =========================
async def _download_telegram_file(bot, file_id: str) -> bytes:
    tg_file = await bot.get_file(file_id)
    bio = tempfile.NamedTemporaryFile(delete=False)
    try:
        await tg_file.download_to_drive(custom_path=bio.name)
        with open(bio.name, "rb") as f:
            return f.read()
    finally:
        try:
            Path(bio.name).unlink(missing_ok=True)
        except Exception:
            pass

def _img_b64(data: bytes) -> str:
    return "data:image/jpeg;base64," + base64.b64encode(data).decode("ascii")

def _summarize_text_with_llm(user_id: int, title: str, text: str) -> str:
    prompt = (
        f"Мне прислали документ «{title}». Сделай короткое резюме и извлеки ключевые пункты.\n\n"
        f"Текст (обрезан до 8000 символов):\n{text[:8000]}"
    )
    return ask_llm(user_id, prompt)

def _analyze_image_with_llm(user_id: int, hint: str, image_b64: str) -> str:
    """
    hint — что хочет пользователь (если пусто — 'опиши что на фото').
    image_b64 — data:image/jpeg;base64,....
    """
    profile = _profile_snapshot(user_id)
    sys_text = TASK_MODES.get(_user_task_mode.get(user_id, "default"), TASK_MODES["default"])["system"]
    instructions = []
    style = profile.get("style")
    language = profile.get("language")
    theme = profile.get("theme")
    if PROFILE_STYLE_INSTRUCTIONS.get(style):
        instructions.append(PROFILE_STYLE_INSTRUCTIONS[style])
    if PROFILE_LANGUAGE_INSTRUCTIONS.get(language):
        instructions.append(PROFILE_LANGUAGE_INSTRUCTIONS[language])
    if PROFILE_THEME_INSTRUCTIONS.get(theme):
        instructions.append(PROFILE_THEME_INSTRUCTIONS[theme])
    if instructions:
        sys_text = f"{sys_text} {' '.join(instructions)}"

    msgs = [
        {"role": "system", "content": sys_text},
        {"role": "user", "content": [
            {"type": "text", "text": hint or "Опиши что на фото и дай ключевые детали."},
            {"type": "image_url", "image_url": {"url": image_b64}},
        ]},
    ]
    return _oai_chat_call(messages=msgs, model="gpt-4o-mini", temperature=0.4)

def _parse_slides_from_text(raw: str, topic: str) -> list[dict[str, list[str]]]:
    slides: list[dict[str, list[str]]] = []
    current_title: str | None = None
    current_bullets: list[str] = []

    for line in (raw or "").splitlines():
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.lower().startswith(("slide", "слайд", "#")) and ":" in stripped:
            if current_title or current_bullets:
                slides.append({
                    "title": current_title or topic,
                    "bullets": current_bullets or ["(нет заметок)"],
                })
            current_title = stripped.split(":", 1)[1].strip() or topic
            current_bullets = []
        else:
            bullet = stripped.lstrip("•*-—– ").strip()
            if bullet:
                current_bullets.append(bullet[:200])

    if current_title or current_bullets:
        slides.append({
            "title": current_title or topic,
            "bullets": current_bullets or ["(нет заметок)"],
        })

    if not slides:
        summary = (raw or "").strip() or "Нет данных"
        slides = [{"title": topic, "bullets": [summary[:200]]}]

    return slides[:8]

_JSON_BLOCK_RE = re.compile(r"```(?:json)?\s*(.*?)\s*```", re.DOTALL | re.IGNORECASE)
_HEX_RE = re.compile(r"^#?([0-9a-fA-F]{6})$")
DEFAULT_PALETTE = {
    "background": "#F5F7FB",
    "accent": "#3F51F9",
    "accent_light": "#E8ECFF",
    "text": "#1F2333",
    "subtitle": "#4D5A7C",
}


def _extract_json_array(raw: str) -> list | None:
    """Пытаемся вытащить JSON-массив из сырого текста."""
    if not raw:
        return None

    def _try_load(candidate: str):
        candidate = candidate.strip()
        if not candidate:
            return None
        try:
            parsed = json.loads(candidate)
            if isinstance(parsed, list):
                return parsed
        except Exception:
            return None
        return None

    direct = _try_load(raw)
    if direct is not None:
        return direct

    for match in _JSON_BLOCK_RE.finditer(raw):
        block = match.group(1)
        parsed = _try_load(block)
        if parsed is not None:
            return parsed

    start = raw.find("[")
    end = raw.rfind("]")
    if start != -1 and end != -1 and end > start:
        parsed = _try_load(raw[start:end + 1])
        if parsed is not None:
            return parsed

    return None


def _normalize_slide(item: dict, idx: int, topic: str) -> dict[str, list[str]]:
    title = str(item.get("title", "")).strip()
    if not title:
        title = f"Раздел {idx}"
    title = title[:120]

    bullets_raw = item.get("bullets")
    if isinstance(bullets_raw, str):
        bullets_iter = [bullets_raw]
    elif isinstance(bullets_raw, list):
        bullets_iter = bullets_raw
    else:
        bullets_iter = []

    bullets: list[str] = []
    for b in bullets_iter:
        if not b:
            continue
        bullet = str(b)
        bullet = re.sub(r"^\s*(?:[-*•]|\d+[.)-])\s*", "", bullet)
        bullet = bullet.strip()
        if not bullet:
            continue
        bullets.append(bullet[:200])

    if not bullets:
        bullets = ["(добавьте заметки сами)"]

    return {"title": title, "bullets": bullets[:5]}


def _extract_json_object(raw: str) -> dict | None:
    """Парсим один JSON-объект из текста."""
    if not raw:
        return None

    def _try(candidate: str):
        candidate = candidate.strip()
        if not candidate:
            return None
        try:
            parsed = json.loads(candidate)
            if isinstance(parsed, dict):
                return parsed
        except Exception:
            return None
        return None

    direct = _try(raw)
    if direct:
        return direct

    for match in _JSON_BLOCK_RE.finditer(raw):
        parsed = _try(match.group(1))
        if parsed:
            return parsed

    start = raw.find("{")
    end = raw.rfind("}")
    if start != -1 and end != -1 and end > start:
        parsed = _try(raw[start:end + 1])
        if parsed:
            return parsed

    return None


def _hex_to_rgb_tuple(value: str | None) -> tuple[int, int, int] | None:
    if not value:
        return None
    match = _HEX_RE.match(value.strip())
    if not match:
        return None
    hex_part = match.group(1)
    return tuple(int(hex_part[i:i + 2], 16) for i in range(0, 6, 2))  # type: ignore[return-value]


def _lighten_color(rgb: tuple[int, int, int], factor: float) -> tuple[int, int, int]:
    factor = max(0.0, min(factor, 1.0))
    return tuple(int(c + (255 - c) * factor) for c in rgb)


def _choose_color_palette(user_id: int, topic: str) -> dict[str, tuple[int, int, int]]:
    prompt = (
        "Подбери гармоничную цветовую палитру для презентации по теме ниже. "
        "Верни JSON-объект без пояснений, формата:\n"
        "{\"background\": \"#RRGGBB\", \"accent\": \"#RRGGBB\", "
        "\"accent_light\": \"#RRGGBB\", \"text\": \"#RRGGBB\", \"subtitle\": \"#RRGGBB\"}\n\n"
        f"Тема: {topic}"
    )
    raw = ask_llm(user_id, prompt)
    data = _extract_json_object(raw) or {}

    palette: dict[str, tuple[int, int, int]] = {}
    for key, fallback in DEFAULT_PALETTE.items():
        rgb = _hex_to_rgb_tuple(str(data.get(key, fallback)))
        if rgb is None:
            rgb = _hex_to_rgb_tuple(fallback)
        palette[key] = rgb or (255, 255, 255)

    if not data.get("accent_light"):
        palette["accent_light"] = _lighten_color(palette["accent"], 0.7)
    return palette


def _pick_slide_emoji(title: str) -> str:
    title_lower = (title or "").lower()
    mapping = [
        ("маркет", "📈"),
        ("продаж", "💼"),
        ("финанс", "💰"),
        ("технолог", "🤖"),
        ("образован", "🎓"),
        ("команда", "🤝"),
        ("анализ", "📊"),
        ("стратег", "🧭"),
        ("дизайн", "🎨"),
        ("риск", "⚠️"),
        ("план", "🗺️"),
        ("экология", "🌱"),
        ("здоров", "🩺"),
        ("продукт", "🧪"),
    ]
    for key, emoji in mapping:
        if key in title_lower:
            return emoji
    return "✨"


def _generate_presentation_structure(user_id: int, topic: str) -> list[dict[str, list[str]]]:
    prompt = (
        "Составь краткую структуру презентации по указанной теме. "
        "Ответь ТОЛЬКО валидным JSON-массивом без префиксов и комментариев. "
        "Формат элемента: {\"title\": \"Название слайда\", \"bullets\": [\"Пункт 1\", \"Пункт 2\"]}. "
        "Нужно 5–7 слайдов: вводный, 3-4 основных, финальный вывод. "
        "Пункты делай короткими (до 15 слов), без номеров и маркеров."
        f"\nТема: {topic!r}"
    )
    raw = ask_llm(user_id, prompt)
    slides_data: list[dict[str, list[str]]] = []

    data = _extract_json_array(raw)
    if data:
        idx = 1
        for item in data:
            if isinstance(item, dict):
                slides_data.append(_normalize_slide(item, idx, topic))
                idx += 1

    if not slides_data:
        parsed = _parse_slides_from_text(raw, topic)
        slides_data = [
            _normalize_slide(item, idx, topic)
            for idx, item in enumerate(parsed, start=1)
        ]

    return slides_data[:8]


async def _generate_presentation_image(topic: str) -> Path | None:
    if not REPLICATE_KEY:
        return None
    prompt = (
        f"High-quality 16:9 illustration for a presentation cover about {topic}. "
        "Modern flat design, soft gradients, no text, professional colour palette."
    )
    try:
        urls = await asyncio.to_thread(_replicate_generate_sync, prompt, width=1280, height=720)
    except Exception:
        return None
    if not urls:
        return None

    try:
        async with httpx.AsyncClient(timeout=30) as client:
            resp = await client.get(urls[0])
            resp.raise_for_status()
            tmpdir = Path(tempfile.gettempdir())
            fpath = tmpdir / f"ppt_cover_{int(time.time())}.png"
            fpath.write_bytes(resp.content)
            return fpath
    except Exception:
        return None


def _set_slide_background(slide, rgb: tuple[int, int, int]):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb)


def _build_presentation_file(
    slides: list[dict[str, list[str]]],
    path: Path,
    topic: str,
    palette: dict[str, tuple[int, int, int]],
    hero_image: Path | None,
):
    prs = Presentation()
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    bg_rgb = palette["background"]
    accent_rgb = palette["accent"]
    accent_light_rgb = palette["accent_light"]
    text_rgb = palette["text"]
    subtitle_rgb = palette["subtitle"]

    title_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(title_slide, bg_rgb)

    top_band = title_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_width, Inches(0.45))
    top_band.fill.solid()
    top_band.fill.fore_color.rgb = RGBColor(*accent_rgb)
    top_band.line.fill.background()

    if hero_image and hero_image.exists():
        try:
            img_width = Inches(4.8)
            img_left = slide_width - img_width - Inches(0.6)
            img_top = Inches(1.1)
            title_slide.shapes.add_picture(str(hero_image), img_left, img_top, width=img_width)
        except Exception:
            pass

    title_box_width = slide_width - Inches(1.5)
    if hero_image and hero_image.exists():
        title_box_width = slide_width - Inches(6.0)
    title_box = title_slide.shapes.add_textbox(Inches(0.8), Inches(1.0), title_box_width, Inches(2.5))
    title_tf = title_box.text_frame
    title_tf.clear()
    title_tf.word_wrap = True
    title_para = title_tf.paragraphs[0]
    title_para.text = topic
    title_para.font.size = Pt(56)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(*text_rgb)

    subtitle_box = title_slide.shapes.add_textbox(Inches(0.8), Inches(3.1), title_box_width, Inches(1))
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.clear()
    subtitle_para = subtitle_tf.paragraphs[0]
    subtitle_para.text = "Сгенерировано NeuroBot 🤖"
    subtitle_para.font.size = Pt(22)
    subtitle_para.font.color.rgb = RGBColor(*subtitle_rgb)

    for idx, slide_data in enumerate(slides, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_slide_background(slide, bg_rgb)

        side_band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), slide_height)
        side_band.fill.solid()
        side_band.fill.fore_color.rgb = RGBColor(*accent_rgb)
        side_band.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), slide_width - Inches(1.6), Inches(1.1))
        title_tf = title_box.text_frame
        title_tf.clear()
        title_para = title_tf.paragraphs[0]
        title_para.text = slide_data.get("title") or f"Слайд {idx}"
        title_para.font.size = Pt(38)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(*accent_rgb)

        emoji_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, slide_width - Inches(1.8), Inches(0.55), Inches(1.05), Inches(1.05))
        emoji_shape.fill.solid()
        emoji_shape.fill.fore_color.rgb = RGBColor(*accent_rgb)
        emoji_shape.line.fill.background()
        emoji_tf = emoji_shape.text_frame
        emoji_tf.clear()
        emoji_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        emoji_para = emoji_tf.paragraphs[0]
        emoji_para.text = _pick_slide_emoji(slide_data.get("title") or "")
        emoji_para.alignment = PP_ALIGN.CENTER
        emoji_para.font.size = Pt(34)
        emoji_para.font.bold = True
        emoji_para.font.color.rgb = RGBColor(*bg_rgb)

        content_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8),
            Inches(1.8),
            slide_width - Inches(1.6),
            slide_height - Inches(2.6),
        )
        content_shape.fill.solid()
        content_shape.fill.fore_color.rgb = RGBColor(*accent_light_rgb)
        content_shape.line.width = Pt(1.8)
        content_shape.line.color.rgb = RGBColor(*accent_rgb)

        text_frame = content_shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.margin_left = Pt(20)
        text_frame.margin_right = Pt(20)
        text_frame.margin_top = Pt(18)
        text_frame.margin_bottom = Pt(18)
        text_frame.vertical_anchor = MSO_ANCHOR.TOP

        bullets = slide_data.get("bullets") or []
        for bullet_idx, bullet in enumerate(bullets):
            para = text_frame.paragraphs[0] if bullet_idx == 0 else text_frame.add_paragraph()
            para.text = bullet
            para.level = 0
            para.font.size = Pt(26 if len(bullets) <= 4 else 22)
            para.font.color.rgb = RGBColor(*text_rgb)
            para.line_spacing = 1.2
            para.space_after = Pt(8)

    prs.save(str(path))

# =========================
# Длинные ответы: нарезка и "Показать ещё"
# =========================
def _split_for_telegram(text: str, limit: int = 3500) -> list[str]:
    """Режет длинные ответы по абзацам, чтобы не рвало середину текста."""
    parts, buf = [], []
    total = 0
    for para in (text or "").split("\n"):
        if total + len(para) + 1 > limit and buf:
            parts.append("\n".join(buf))
            buf, total = [], 0
        buf.append(para)
        total += len(para) + 1
    if buf:
        parts.append("\n".join(buf))
    return parts if parts else [text]

async def on_more_btn(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Отправляет следующую часть длинного ответа и даёт кнопки 'Показать ещё' + 'Озвучить'."""
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass

    uid = q.from_user.id
    queue = _long_reply_queue.get(uid) or []
    if not queue:
        # нечего слать — просто уберём клавиатуру
        try:
            await q.message.edit_reply_markup(reply_markup=None)
        except Exception:
            pass
        return

    next_part = queue.pop(0)
    _long_reply_queue[uid] = queue

    # важно: озвучиваем именно тот кусок, который сейчас показываем
    _last_answer[uid] = next_part

    # если ещё есть части — две кнопки, иначе оставим только «Озвучить»
    rows: list[list[InlineKeyboardButton]] = []
    if queue:
        rows.append([
            InlineKeyboardButton("Показать ещё ▶️", callback_data="more"),
            InlineKeyboardButton("🎧 Озвучить", callback_data="tts"),
        ])
    else:
        rows.append([InlineKeyboardButton("🎧 Озвучить", callback_data="tts")])
    rows.append([
        InlineKeyboardButton("⭐ Шаблон", callback_data="fav:add"),
        InlineKeyboardButton("🔁 Перевести", callback_data="quick:translate"),
        InlineKeyboardButton("🧾 Сжать", callback_data="quick:summary"),
    ])

    await q.message.reply_text(next_part, reply_markup=InlineKeyboardMarkup(rows))

# ---------- Images (Replicate: Flux-1 Schnell) ----------
def _replicate_generate_sync(prompt: str, width: int = 1024, height: int = 1024) -> list[str]:
    """
    Блокирующая генерация через Replicate. Возвращает список URL готовых изображений.
    """
    if not REPLICATE_KEY:
        raise RuntimeError("REPLICATE_KEY пуст — подключите ключ Replicate в .env")

    model = "black-forest-labs/flux-schnell"
    headers = {
        "Authorization": f"Token {REPLICATE_KEY}",
        "Content-Type": "application/json"
    }

    create_payload = {
        "input": {
            "prompt": prompt,
            "width": width,
            "height": height,
            "num_outputs": 1,
            "go_fast": True
        }
    }

    create = requests.post(
        f"https://api.replicate.com/v1/models/{model}/predictions",
        json=create_payload,
        headers=headers,
        timeout=30
    )
    create.raise_for_status()
    prediction = create.json()
    pred_id = prediction.get("id")
    if not pred_id:
        raise RuntimeError(f"Replicate: не получили id предсказания: {prediction}")

    status = prediction.get("status")
    get_url = f"https://api.replicate.com/v1/predictions/{pred_id}"

    for _ in range(60):
        if status in ("succeeded", "failed", "canceled"):
            break
        poll = requests.get(get_url, headers=headers, timeout=15)
        poll.raise_for_status()
        prediction = poll.json()
        status = prediction.get("status")
        if status == "succeeded":
            break
        time.sleep(1)

    if status != "succeeded":
        err = prediction.get("error") or status
        raise RuntimeError(f"Replicate: задача не удалась: {err}")

    output = prediction.get("output") or []
    if isinstance(output, str):
        output = [output]
    return output


_CYRILLIC_RE = re.compile(r"[А-Яа-яЁё]")


def _translate_to_english(text: str) -> str:
    try:
        msgs = [
            {
                "role": "system",
                "content": (
                    "You are a professional translator. Translate the user's prompt into concise, natural English "
                    "suitable for an image generation model. Respond with the translation only."
                ),
            },
            {"role": "user", "content": text},
        ]
        translated = _oai_chat_call(messages=msgs, model="gpt-4o-mini", temperature=0.0)
        return translated.strip()
    except Exception:
        return text


def _prepare_image_prompt(prompt: str) -> str:
    text = (prompt or "").strip()
    if not text:
        return text

    if _CYRILLIC_RE.search(text):
        translated = _translate_to_english(text)
        translated = translated.strip().strip('"').strip("'")
        if translated and translated.lower() != "none":
            return f"{translated}. Original description (Russian): {text}"
    return text


async def generate_image_and_send(user_id: int, chat_id: int, prompt: str, bot) -> None:
    try:
        prepared_prompt = _prepare_image_prompt(prompt)
        urls = await asyncio.to_thread(_replicate_generate_sync, prepared_prompt)
        if not urls:
            await bot.send_message(chat_id=chat_id, text="Не удалось получить изображение.")
            return
        await bot.send_photo(chat_id=chat_id, photo=urls[0], caption="Готово ✅")
    except Exception as e:
        await bot.send_message(chat_id=chat_id, text=f"Ошибка генерации: {e}")

# ---------- Favorites & Быстрые действия ----------

def _short_title(text: str) -> str:
    cleaned = " ".join(text.strip().split())
    if not cleaned:
        return "Без названия"
    return textwrap.shorten(cleaned, width=32, placeholder="…")


async def _favorites_payload(user_id: int) -> tuple[str, InlineKeyboardMarkup]:
    favs = await list_favorite_prompts(user_id)
    if not favs:
        text = (
            "⭐ <b>Шаблоны</b>\n"
            "У вас пока нет сохранённых подсказок.\n"
            "Нажмите кнопку «⭐ Шаблон» под ответом, чтобы добавить текущий запрос в избранное."
        )
        kb = InlineKeyboardMarkup([[InlineKeyboardButton("⬅️ Назад", callback_data="home")]])
        return text, kb

    lines = ["⭐ <b>Шаблоны</b>\nВыберите действие:"]
    rows: list[list[InlineKeyboardButton]] = []
    for fid, title in favs[:10]:
        short = _short_title(title)
        rows.append([
            InlineKeyboardButton(f"▶️ {short}", callback_data=f"fav:run:{fid}"),
            InlineKeyboardButton("🗑️", callback_data=f"fav:del:{fid}"),
        ])
    rows.append([InlineKeyboardButton("⬅️ Назад", callback_data="home")])
    text = "\n".join(lines)
    return text, InlineKeyboardMarkup(rows)


async def cmd_favorites(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    text, kb = await _favorites_payload(user_id)
    await update.message.reply_text(text, parse_mode="HTML", reply_markup=kb)


async def on_favorites_btn(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass
    text, kb = await _favorites_payload(q.from_user.id)
    try:
        await q.message.edit_text(text, parse_mode="HTML", reply_markup=kb)
    except Exception:
        await q.message.reply_text(text, parse_mode="HTML", reply_markup=kb)


async def on_fav_add(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass
    uid = q.from_user.id
    prompt = _last_user_prompt.get(uid)
    if not prompt:
        await q.message.reply_text("Нет последнего запроса, нечего сохранить.", reply_markup=main_keyboard())
        return
    title = _short_title(prompt)
    fav_id = await add_favorite_prompt(uid, title, prompt)
    logger.info("Saved favorite prompt %s for %s", fav_id, uid)
    await q.message.reply_text(f"⭐ Шаблон «{title}» сохранён. Откройте меню «⭐ Шаблоны», чтобы использовать его.", reply_markup=main_keyboard())


async def on_fav_delete(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass
    uid = q.from_user.id
    fav_id = int(q.data.split("fav:del:", 1)[-1])
    ok = await delete_favorite_prompt(uid, fav_id)
    if ok:
        text, kb = await _favorites_payload(uid)
        try:
            await q.message.edit_text(text, parse_mode="HTML", reply_markup=kb)
        except Exception:
            await q.message.reply_text("Шаблон удалён.", reply_markup=main_keyboard())
    else:
        await q.message.reply_text("Не удалось удалить шаблон.", reply_markup=main_keyboard())


async def on_fav_run(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    try:
        await q.answer("Использую шаблон…")
    except Exception:
        pass
    uid = q.from_user.id
    fav_id = int(q.data.split("fav:run:", 1)[-1])
    fav = await get_favorite_prompt(uid, fav_id)
    if not fav:
        await q.message.reply_text("Шаблон не найден.", reply_markup=main_keyboard())
        return
    _, prompt_text = fav
    fake_update = SimpleNamespace(message=q.message, effective_user=q.from_user)
    await _handle_text_request(fake_update, context, prompt_text)


def _detect_translation_target(profile: dict[str, str], text: str) -> tuple[str, str]:
    pref_lang = profile.get("language", "auto")
    if pref_lang == "ru":
        return "English", "английский язык"
    if pref_lang == "en":
        return "Russian", "русский язык"
    if _CYRILLIC_RE.search(text):
        return "English", "английский язык"
    return "Russian", "русский язык"


async def on_quick_translate(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass
    uid = q.from_user.id
    text = _last_answer.get(uid)
    if not text:
        await q.message.reply_text("Нет текста для перевода.", reply_markup=main_keyboard())
        return
    await _ensure_profile(uid)
    profile = _profile_snapshot(uid)
    target_code, target_label = _detect_translation_target(profile, text)
    prompt = (
        f"Translate the text below into {target_code}. Respond with the translation only.\n\n{text}"
        if target_code == "English"
        else f"Переведи текст ниже на {target_label}. Передай только перевод без комментариев.\n\n{text}"
    )
    try:
        translation = _oai_chat_call(
            messages=[
                {"role": "system", "content": "You are a precise translator."},
                {"role": "user", "content": prompt},
            ],
            model="gpt-4o-mini",
            temperature=0,
        ).strip()
    except Exception as e:
        await q.message.reply_text(f"Не удалось перевести: {e}", reply_markup=main_keyboard())
        return
    label = "Перевод (EN)" if target_code == "English" else "Перевод (RU)"
    await q.message.reply_text(f"{label}:\n{translation}", reply_markup=main_keyboard())


async def on_quick_summary(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass
    uid = q.from_user.id
    text = _last_answer.get(uid)
    if not text:
        await q.message.reply_text("Нет текста для сжатия.", reply_markup=main_keyboard())
        return
    await _ensure_profile(uid)
    profile = _profile_snapshot(uid)
    lang = profile.get("language", "auto")
    if lang == "en":
        prompt = f"Summarize the text below in 3-4 bullet points.\n\n{text}"
    else:
        prompt = f"Сделай краткое резюме текста ниже в 3–4 пунктах.\n\n{text}"
    try:
        summary = _oai_chat_call(
            messages=[
                {"role": "system", "content": "You create short helpful summaries."},
                {"role": "user", "content": prompt},
            ],
            model="gpt-4o-mini",
            temperature=0.2,
        ).strip()
    except Exception as e:
        await q.message.reply_text(f"Не удалось создать резюме: {e}", reply_markup=main_keyboard())
        return
    await q.message.reply_text(f"🧾 Краткое резюме:\n{summary}", reply_markup=main_keyboard())

# ---------- Настройки профиля ----------

def _settings_text(profile: dict[str, str]) -> str:
    return (
        "⚙️ <b>Персональные настройки</b>\n\n"
        f"• Стиль: <b>{PROFILE_STYLES.get(profile.get('style'), 'Стандарт')}</b>\n"
        f"• Язык: <b>{PROFILE_LANGUAGES.get(profile.get('language'), 'Авто')}</b>\n"
        f"• Формат: <b>{PROFILE_FORMATS.get(profile.get('output_format'), 'Обычный')}</b>\n"
        f"• Тема карточек: <b>{PROFILE_THEMES.get(profile.get('theme'), 'Авто')}</b>\n\n"
        "Выберите параметр, чтобы изменить его."
    )


def _settings_keyboard(profile: dict[str, str]) -> InlineKeyboardMarkup:
    style_buttons = [
        InlineKeyboardButton(
            ("✅ " if profile.get("style") == key else "") + label,
            callback_data=f"settings:style:{key}"
        )
        for key, label in PROFILE_STYLES.items()
    ]
    language_buttons = [
        InlineKeyboardButton(
            ("✅ " if profile.get("language") == key else "") + label,
            callback_data=f"settings:language:{key}"
        )
        for key, label in PROFILE_LANGUAGES.items()
    ]
    format_buttons = [
        InlineKeyboardButton(
            ("✅ " if profile.get("output_format") == key else "") + label,
            callback_data=f"settings:format:{key}"
        )
        for key, label in PROFILE_FORMATS.items()
    ]
    theme_buttons = [
        InlineKeyboardButton(
            ("✅ " if profile.get("theme") == key else "") + label,
            callback_data=f"settings:theme:{key}"
        )
        for key, label in PROFILE_THEMES.items()
    ]

    def chunk(buttons: list[InlineKeyboardButton], size: int = 3) -> list[list[InlineKeyboardButton]]:
        return [buttons[i:i + size] for i in range(0, len(buttons), size)]

    rows: list[list[InlineKeyboardButton]] = []
    rows.extend(chunk(style_buttons, size=2))
    rows.extend(chunk(language_buttons, size=3))
    rows.extend(chunk(format_buttons, size=3))
    rows.extend(chunk(theme_buttons, size=3))
    rows.append([InlineKeyboardButton("⬅️ Назад", callback_data="home")])
    return InlineKeyboardMarkup(rows)


async def cmd_settings(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    profile = await _ensure_profile(user_id)
    await update.message.reply_text(
        _settings_text(profile),
        parse_mode="HTML",
        reply_markup=_settings_keyboard(profile),
    )


async def on_settings_btn(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass
    profile = await _ensure_profile(q.from_user.id)
    try:
        await q.message.edit_text(
            _settings_text(profile),
            parse_mode="HTML",
            reply_markup=_settings_keyboard(profile),
        )
    except Exception:
        await q.message.reply_text(
            _settings_text(profile),
            parse_mode="HTML",
            reply_markup=_settings_keyboard(profile),
        )


async def on_settings_change(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass
    try:
        _, field, value = q.data.split(":", 2)
    except ValueError:
        return
    user_id = q.from_user.id

    allowed = {
        "style": set(PROFILE_STYLES.keys()),
        "language": set(PROFILE_LANGUAGES.keys()),
        "format": set(PROFILE_FORMATS.keys()),
        "theme": set(PROFILE_THEMES.keys()),
    }
    if field not in allowed or value not in allowed[field]:
        await q.message.reply_text("Некорректное значение.", reply_markup=main_keyboard())
        return

    await set_user_profile_value(user_id, field, value)
    _update_profile_cache(user_id, field, value)
    profile = await _ensure_profile(user_id)
    try:
        await q.message.edit_text(
            _settings_text(profile),
            parse_mode="HTML",
            reply_markup=_settings_keyboard(profile),
        )
    except Exception:
        await q.message.reply_text(
            _settings_text(profile),
            parse_mode="HTML",
            reply_markup=_settings_keyboard(profile),
        )
# ---------- UI ----------
def main_keyboard() -> InlineKeyboardMarkup:
    return InlineKeyboardMarkup([
        [InlineKeyboardButton("🧠 Выбрать модель", callback_data="models")],
        [InlineKeyboardButton("🎛 Режимы", callback_data="modes")],
        [InlineKeyboardButton("💬 Диалоги", callback_data="dialog")],
        [InlineKeyboardButton("🖼️ Создать картинку", callback_data="img"),
         InlineKeyboardButton("🗂️ Презентация", callback_data="ppt")],
        [InlineKeyboardButton("👤 Профиль", callback_data="profile"),
         InlineKeyboardButton("⚙️ Настройки", callback_data="settings")],
        [InlineKeyboardButton("⭐ Шаблоны", callback_data="fav:list"),
         InlineKeyboardButton("🎁 Реферальная программа", callback_data="ref")],
        [InlineKeyboardButton("❓ Помощь", callback_data="help:how"),
         InlineKeyboardButton("📚 FAQ",    callback_data="help:faq")],
        [InlineKeyboardButton("💳 Купить подписку", callback_data="buy")],
    ])

# ===== Меню моделей =====
def _models_menu_text(mode: str = "short") -> str:
    if mode == "short":
        return (
            "<b>Кратко о моделях</b>\n"
            "• <b>GPT-5</b> — флагман для сложных задач, кодинга и длинных контекстов.\n"
            "• <b>Claude 4.5 Sonnet</b> — силён в анализе, стиле и длинных ответах.\n"
            "• <b>Gemini 2.5 Pro</b> — хороший баланс скорости и качества, мультимодальность.\n"
            "• <b>OpenAI o3</b> — логика и рассуждения, аккуратный тон.\n"
            "• <b>DeepSeek V3.2</b> — быстрые и экономные ответы, отлично для повседневки.\n"
            "• <b>OpenAI o4-mini</b> — быстрые короткие ответы и прототипирование.\n"
            "• <b>GPT-5 mini</b> — лёгкая версия для черновиков и быстрых итераций.\n"
            "• <b>GPT-4o search</b> — модель с упором на поиск/извлечение фактов.\n"
            "• <b>GPT-4o mini</b> — экономичная альтернатива для простых задач.\n"
            "• <b>Claude 3.5 Haiku</b> — очень быстро на коротких запросах.\n"
            "• <b>Gemini 2.5 Flash</b> — быстрые черновики, резюме, списки.\n\n"
            "Выберите модель для работы:"
        )
    else:
        return (
            "<b>Подробно о моделях</b>\n"
            "<b>GPT-5</b> — топ по качеству кода, сложным рассуждениям и длинным контекстам. Рекомендуется для архитектуры, аудитов, сложных SQL и интеграций.\n\n"
            "<b>Claude 4.5 Sonnet</b> — силён в языке и стиле: эссе, стратегии, юридические/деловые тексты, аккуратные объяснения. Хорош на огромных документах.\n\n"
            "<b>Gemini 2.5 Pro</b> — сбалансирован: анализ, идеи, мультимодальные задачи. Подходит для презентаций, маркетинга и быстрых исследований.\n\n"
            "<b>OpenAI o3</b> — фокус на логике/Chain-of-Thought: пошаговые решения, математика, тонкая аргументация, проверка гипотез.\n\n"
            "<b>DeepSeek V3.2</b> — очень быстрый и экономичный: повседневные вопросы, резюме, генерация простых текстов и черновики.\n\n"
            "<b>OpenAI o4-mini</b> — быстрые ответы и прототипирование: черновые спецификации, user stories, наброски кода.\n\n"
            "<b>GPT-5 mini</b> — минимальные задержки: идеи, списки, короткие подсказки, быстрые итерации.\n\n"
            "<b>GPT-4o search</b> — приоритизирует поиск и извлечение: набор фактов, цитаты, обзорные справки.\n\n"
            "<b>GPT-4o mini</b> — экономичный универсал для простых задач/переводов и быстрых советов.\n\n"
            "<b>Claude 3.5 Haiku</b> — молниеносные короткие ответы и рефакторинг текста, подсветка смысла.\n\n"
            "<b>Gemini 2.5 Flash</b> — резюме страниц, TODO-списки, короткие письма, быстрые описания.\n\n"
            "Выберите модель:"
        )

def models_keyboard_visual() -> InlineKeyboardMarkup:
    return InlineKeyboardMarkup([
        [InlineKeyboardButton("🔸 Кратко",  callback_data="mvis:short"),
         InlineKeyboardButton("ℹ️ Подробно", callback_data="mvis:full")],
        [InlineKeyboardButton("Claude 3.5 Haiku", callback_data="mvis:sel:Claude 3.5 Haiku"),
         InlineKeyboardButton("✅ GPT-5",         callback_data="mvis:sel:GPT-5")],
        [InlineKeyboardButton("Claude 4.5 Sonnet", callback_data="mvis:sel:Claude 4.5 Sonnet"),
         InlineKeyboardButton("Gemini 2.5 Pro",    callback_data="mvis:sel:Gemini 2.5 Pro")],
        [InlineKeyboardButton("OpenAI o3",         callback_data="mvis:sel:OpenAI o3"),
         InlineKeyboardButton("DeepSeek V3.2",     callback_data="mvis:sel:DeepSeek V3.2")],
        [InlineKeyboardButton("OpenAI o4-mini",    callback_data="mvis:sel:OpenAI o4-mini"),
         InlineKeyboardButton("GPT-5 mini",        callback_data="mvis:sel:GPT-5 mini")],
        [InlineKeyboardButton("GPT-4o search 🔎",  callback_data="mvis:sel:GPT-4o search"),
         InlineKeyboardButton("GPT-4o mini",       callback_data="mvis:sel:GPT-4o mini")],
        [InlineKeyboardButton("Gemini 2.5 Flash",  callback_data="mvis:sel:Gemini 2.5 Flash")],
        [InlineKeyboardButton("⬅️ Назад",          callback_data="home")],
    ])

# ===== Меню режимов (ярлыки) =====
def modes_keyboard() -> InlineKeyboardMarkup:
    return InlineKeyboardMarkup([
        [InlineKeyboardButton("Стандарт", callback_data="mode:default"),
         InlineKeyboardButton("Кодинг",   callback_data="mode:coding")],
        [InlineKeyboardButton("SEO",      callback_data="mode:seo"),
         InlineKeyboardButton("Перевод",  callback_data="mode:translate")],
        [InlineKeyboardButton("Резюме",   callback_data="mode:summarize"),
         InlineKeyboardButton("Креатив",  callback_data="mode:creative")],
        [InlineKeyboardButton("⬅️ Назад", callback_data="home")],
    ])

def current_mode_label(user_id: int) -> str:
    key = _user_task_mode.get(user_id, "default")
    return TASK_MODES.get(key, TASK_MODES["default"])["label"]

# ===== Диалоговые режимы (simple / rooms) =====

def dialog_menu_text(mode: str) -> str:
    common_note = "\n\n<i>ℹ️ Контекстный режим работает с любой выбранной моделью.</i>"
    if mode == DIALOG_ROOMS:
        return (
            "<b>Диалоги с контекстом</b>\n"
            "Создавайте отдельные чаты по темам: история сообщений сохраняется и учитывается в ответах."
            f"{common_note}"
        )
    else:
        return (
            "<b>Быстрые ответы</b>\n"
            "Каждое сообщение — независимое. История не копится, ответы максимально быстрые."
            f"{common_note}"
        )

def dialog_keyboard(mode_now: str) -> InlineKeyboardMarkup:
    kb = [
        [InlineKeyboardButton(
            ("✅ " if mode_now == DIALOG_SIMPLE else "") + "⚡ Быстрые ответы",
            callback_data="dialog:simple"
        )],
        [InlineKeyboardButton(
            ("✅ " if mode_now == DIALOG_ROOMS else "") + "🗂️ Диалоги с контекстом",
            callback_data="dialog:rooms"
        )],
        [InlineKeyboardButton("📂 Мои чаты", callback_data="chats")],
        [InlineKeyboardButton("⬅️ Назад", callback_data="home")],
    ]
    return InlineKeyboardMarkup(kb)

async def on_dialog_btn(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass
    mode = await get_chat_mode(q.from_user.id)
    await q.message.edit_text(dialog_menu_text(mode), parse_mode="HTML", reply_markup=dialog_keyboard(mode))

async def on_dialog_select(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass
    want = q.data.split("dialog:", 1)[-1]
    want = DIALOG_ROOMS if want == "rooms" else DIALOG_SIMPLE
    await set_chat_mode(q.from_user.id, want)

    # Если включили rooms и нет активного чата — создадим первый
    if want == DIALOG_ROOMS:
        active = await get_active_chat(q.from_user.id)
        if active is None:
            cid = await create_chat(q.from_user.id, "Чат 1")
            await set_active_chat(q.from_user.id, cid)

    await q.message.edit_text(dialog_menu_text(want), parse_mode="HTML", reply_markup=dialog_keyboard(want))

async def on_chats_btn(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass
    user_id = q.from_user.id
    await set_chat_mode(user_id, DIALOG_ROOMS)  # при открытии списка чатов — сразу режим rooms
    chats = await list_chats(user_id)
    active = await get_active_chat(user_id)

    rows = []
    if not chats:
        rows.append([InlineKeyboardButton("➕ Создать первый чат", callback_data="chat:new")])
    else:
        for cid, title, pinned in chats[:10]:
            prefix = "✅ " if active == cid else ""
            if pinned:
                prefix = f"{prefix}📌 "
            rows.append([InlineKeyboardButton(f"{prefix}{title}", callback_data=f"chat:open:{cid}")])
        rows.append([InlineKeyboardButton("➕ Новый чат", callback_data="chat:new")])

    rows.append([InlineKeyboardButton("⬅️ Назад", callback_data="dialog")])
    await q.message.edit_text("Ваши чаты:", reply_markup=InlineKeyboardMarkup(rows))

async def on_chat_new(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass
    user_id = q.from_user.id
    chats = await list_chats(user_id)
    title = f"Чат {len(chats)+1}"
    cid = await create_chat(user_id, title)
    await set_active_chat(user_id, cid)
    await on_chats_btn(update, context)

async def on_chat_open(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass
    user_id = q.from_user.id
    cid = int(q.data.split("chat:open:", 1)[-1])
    await set_active_chat(user_id, cid)

    # найдём заголовок чата
    chats = await list_chats(user_id)
    title = next((t for (i, t, _) in chats if i == cid), f"Чат {cid}")
    pinned = next((p for (i, _, p) in chats if i == cid), False)

    pin_label = "📌 Закрепить" if not pinned else "📍 Открепить"
    kb = InlineKeyboardMarkup([
        [InlineKeyboardButton("✏️ Переименовать", callback_data=f"chat:rename:{cid}")],
        [InlineKeyboardButton(pin_label, callback_data=f"chat:pin:{cid}")],
        [InlineKeyboardButton("🔗 Поделиться ссылкой", callback_data=f"chat:share:{cid}")],
        [InlineKeyboardButton("📄 Экспорт PDF", callback_data=f"chat:export:pdf:{cid}")],
        [InlineKeyboardButton("🧾 Markdown для Notion", callback_data=f"chat:export:md:{cid}")],
        [InlineKeyboardButton("🗑️ Удалить",       callback_data=f"chat:delete:{cid}")],
        [InlineKeyboardButton("⬅️ К списку чатов", callback_data="chats")]
    ])
    status_line = "📌 Закреплён" if pinned else "📎 Не закреплён"
    await q.message.edit_text(
        f"Чат: <b>{title}</b>\n{status_line}\nВыберите действие:",
        parse_mode="HTML",
        reply_markup=kb,
    )

async def on_chat_rename_ask(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass
    user_id = q.from_user.id
    cid = int(q.data.split("chat:rename:", 1)[-1])
    _pending_chat_rename[user_id] = cid
    kb = InlineKeyboardMarkup([[InlineKeyboardButton("⬅️ Отмена", callback_data="chats")]])
    await q.message.edit_text("Отправьте новое название чата (1–80 символов):", reply_markup=kb)

async def on_chat_delete_confirm(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass
    cid = int(q.data.split("chat:delete:", 1)[-1])
    kb = InlineKeyboardMarkup([
        [InlineKeyboardButton("✅ Да, удалить", callback_data=f"chat:delete:do:{cid}")],
        [InlineKeyboardButton("⬅️ Отмена", callback_data="chats")]
    ])
    await q.message.edit_text("Удалить этот чат? Действие необратимо.", reply_markup=kb)

async def on_chat_delete_do(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass
    user_id = q.from_user.id
    cid = int(q.data.split("chat:delete:do:", 1)[-1])

    # если удаляем активный — потом сбросим active_chat_id
    active = await get_active_chat(user_id)
    ok = await delete_chat(user_id, cid)
    if ok and active == cid:
        await set_active_chat(user_id, None)

    # если после удаления нет чатов — предложим создать первый
    chats = await list_chats(user_id)
    if not chats:
        kb = InlineKeyboardMarkup([
            [InlineKeyboardButton("➕ Создать чат", callback_data="chat:new")],
            [InlineKeyboardButton("⬅️ Назад", callback_data="dialog")]
        ])
        await q.message.edit_text("Чат удалён. У вас пока нет чатов.", reply_markup=kb)
        return

    # иначе вернёмся к списку
    await on_chats_btn(update, context)


async def _get_chat_meta(user_id: int, chat_id: int) -> tuple[str, bool]:
    chats = await list_chats(user_id)
    for cid, title, pinned in chats:
        if cid == chat_id:
            return title, pinned
    return f"Чат {chat_id}", False


async def on_chat_pin_toggle(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    user_id = q.from_user.id
    chat_id = int(q.data.split("chat:pin:", 1)[-1])
    title, pinned = await _get_chat_meta(user_id, chat_id)
    await set_chat_pinned(user_id, chat_id, not pinned)
    try:
        await q.answer(f"Чат «{title}» {'закреплён' if not pinned else 'откреплён'}.", show_alert=False)
    except Exception:
        pass
    await on_chat_open(update, context)


def _chat_history_to_markdown(title: str, history: list[tuple[str, str, str]]) -> str:
    lines = [f"# {title}", ""]
    for role, content, created_at in history:
        label = "Пользователь" if role == "user" else "Ассистент" if role == "assistant" else "Система"
        timestamp = created_at.replace("T", " ").split("+", 1)[0]
        lines.append(f"## {label} · {timestamp}")
        lines.append(content.strip())
        lines.append("")
    return "\n".join(lines)


def _chat_history_to_html(title: str, history: list[tuple[str, str, str]]) -> str:
    body = [f"<h1>{html.escape(title)}</h1>"]
    for role, content, created_at in history:
        label = "Пользователь" if role == "user" else "Ассистент" if role == "assistant" else "Система"
        timestamp = created_at.replace("T", " ").split("+", 1)[0]
        body.append("<div class='entry'>")
        body.append(f"<div class='meta'>{html.escape(label)} · {html.escape(timestamp)}</div>")
        safe = html.escape(content).replace("\n", "<br>")
        body.append(f"<div class='content'>{safe}</div>")
        body.append("</div>")
    return "\n".join(body)


def _find_font_path() -> str | None:
    candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/dejavu/DejaVuSans.ttf",
        "/Library/Fonts/Arial Unicode.ttf",
        "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
        str(Path.home() / "Library/Fonts/Arial Unicode.ttf"),
    ]
    for path in candidates:
        if Path(path).exists():
            return path
    return None


def _build_pdf_from_history(title: str, history: list[tuple[str, str, str]], dest: Path):
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    font_path = _find_font_path()
    if font_path:
        pdf.add_font("Custom", "", font_path, uni=True)
        font_name = "Custom"
    else:
        font_name = "Arial"
    pdf.add_page()
    pdf.set_font(font_name, size=18)
    pdf.multi_cell(0, 10, title)
    pdf.ln(4)
    for role, content, created_at in history:
        label = "Пользователь" if role == "user" else "Ассистент" if role == "assistant" else "Система"
        timestamp = created_at.replace("T", " ").split("+", 1)[0]
        pdf.set_font(font_name, size=12)
        header = f"{label} · {timestamp}"
        if not font_path:
            header = header.encode("latin-1", "replace").decode("latin-1")
        pdf.multi_cell(0, 7, header)
        pdf.set_font(font_name, size=11)
        text = content.strip()
        if not font_path:
            text = text.encode("latin-1", "replace").decode("latin-1")
        pdf.multi_cell(0, 6, text)
        pdf.ln(4)
    pdf.output(str(dest))


async def on_chat_share(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass
    user_id = q.from_user.id
    chat_id = int(q.data.split("chat:share:", 1)[-1])
    if not _public_url:
        await q.message.reply_text("Общая ссылка недоступна: не задан PUBLIC_URL.", reply_markup=main_keyboard())
        return
    title, _ = await _get_chat_meta(user_id, chat_id)
    await cleanup_chat_shares()
    token, expires_iso = await create_chat_share(user_id, chat_id)
    link = f"{_public_url.rstrip('/')}/share/{token}"
    expires_dt = datetime.fromisoformat(expires_iso)
    expires_text = expires_dt.strftime("%d.%m.%Y %H:%M")
    await q.message.reply_text(
        f"🔗 Ссылка на чат «{title}»:\n{link}\n\n"
        f"Действует до {expires_text} (UTC). Передайте ссылку, чтобы коллеги могли просмотреть диалог.",
        disable_web_page_preview=True,
        reply_markup=main_keyboard(),
    )


async def on_chat_export_md(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass
    user_id = q.from_user.id
    chat_id = int(q.data.split("chat:export:md:", 1)[-1])
    title, _ = await _get_chat_meta(user_id, chat_id)
    history = await get_chat_history_all(chat_id)
    if not history:
        await q.message.reply_text("Чат пуст — нечего экспортировать.", reply_markup=main_keyboard())
        return
    markdown = _chat_history_to_markdown(title, history)
    tmpdir = Path(tempfile.gettempdir())
    fname = re.sub(r"[^A-Za-z0-9]+", "_", title)[:40] or f"chat_{chat_id}"
    path = tmpdir / f"{fname}.md"
    path.write_text(markdown, encoding="utf-8")
    try:
        with open(path, "rb") as fh:
            await q.message.reply_document(
                document=fh,
                filename=path.name,
                caption="Экспортирован в Markdown. Импортируйте файл в Notion или откройте в редакторе.",
            )
    finally:
        try:
            path.unlink(missing_ok=True)
        except Exception:
            pass


async def on_chat_export_pdf(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass
    user_id = q.from_user.id
    chat_id = int(q.data.split("chat:export:pdf:", 1)[-1])
    title, _ = await _get_chat_meta(user_id, chat_id)
    history = await get_chat_history_all(chat_id)
    if not history:
        await q.message.reply_text("Чат пуст — нечего экспортировать.", reply_markup=main_keyboard())
        return
    tmpdir = Path(tempfile.gettempdir())
    fname = re.sub(r"[^A-Za-z0-9]+", "_", title)[:40] or f"chat_{chat_id}"
    path = tmpdir / f"{fname}.pdf"
    try:
        _build_pdf_from_history(title, history, path)
        with open(path, "rb") as fh:
            await q.message.reply_document(
                document=fh,
                filename=path.name,
                caption="PDF-файл готов — можно делиться с командой или печатать.",
            )
    except Exception as e:
        await q.message.reply_text(f"Не удалось собрать PDF: {e}", reply_markup=main_keyboard())
    finally:
        try:
            path.unlink(missing_ok=True)
        except Exception:
            pass

# =========================
# Кнопка/команда генерации изображений
# =========================
async def on_img_btn(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass

    user_id = q.from_user.id
    if not await is_premium(user_id):
        await q.message.reply_text(
            "Доступно в Премиум.\n\n"
            "Премиум даёт:\n"
            "• Безлимитные сообщения\n"
            "• Доступ ко всем моделям\n"
            "• Генерацию изображений\n\n"
            f"Нажмите «Купить подписку», стоимость {PRICE_RUB_TEXT} на 30 дней.",
            reply_markup=InlineKeyboardMarkup(
                [[InlineKeyboardButton("Купить подписку", callback_data="buy")]]
            )
        )
        return

    _awaiting_img_prompt[user_id] = True
    await q.message.reply_text("Опиши картинку текстом (1–2 предложения). Я сгенерирую изображение.")

async def cmd_img(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    if not await is_premium(user_id):
        await update.message.reply_text(
            f"Генерация изображений доступна в Премиум ({PRICE_RUB_TEXT} / 30 дней).",
            reply_markup=InlineKeyboardMarkup(
                [[InlineKeyboardButton("Купить подписку", callback_data="buy")]]
            )
        )
        return
    _awaiting_img_prompt[user_id] = True
    await update.message.reply_text(
        "Опиши картинку текстом. Пример: «синий неоновый город, дождь, стиль киберпанк»."
    )

# =========================
# /start + рефералка
# =========================
async def cmd_start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    await add_user(user.id)

    # deep-link параметр: /start ref_<tg_id>
    ref_id = None
    if context.args:
        arg = context.args[0]
        if arg.startswith("ref_"):
            try:
                ref_id = int(arg.split("ref_", 1)[1])
            except Exception:
                ref_id = None

    if ref_id and ref_id != user.id:
        try:
            first_bind = await set_referrer_if_empty(user.id, ref_id)
            if first_bind:
                await add_free_credits(ref_id, REF_BONUS)
                try:
                    await application.bot.send_message(
                        chat_id=ref_id,
                        text=f"🎉 По вашей ссылке зарегистрировался новый пользователь.\n+{REF_BONUS} бесплатных заявок!"
                    )
                except Exception:
                    pass
        except Exception as e:
            logger.warning("ref attach failed: %s", e)

    # НОВОЕ ПРИВЕТСТВИЕ
    text = (
        "Привет! Я <b>НейроБот 🤖</b> — твой умный помощник для текста, идей, кода и перевода.\n\n"
        f"🆓 <b>Бесплатный доступ</b> — {DAILY_LIMIT} сообщений в день + бонусы за приглашённых друзей.\n"
        "💎 <b>Премиум</b> — без ограничений, очередей и лимитов, доступ ко всем моделям и генерации изображений. "
        f"Стоимость — всего <b>{PRICE_RUB_TEXT} на 30 дней</b>.\n\n"
        "🚀 Что я умею:\n"
        "• Отвечать на вопросы и писать тексты любой сложности\n"
        "• Помогать с кодом и объяснять ошибки\n"
        "• Переводить 🇷🇺↔️🇬🇧 тексты и документы\n"
        "• Генерировать идеи, резюме, описания и письма\n"
        "• Создавать изображения по описанию 🖼️\n"
        "• 🎧 <b>Озвучивать ответы голосом</b> — нажми кнопку «Озвучить» под сообщением\n"
        "• 📄 <b>Работать с документами</b> (.txt, .md, .csv, .pdf): краткие выжимки и ключевые пункты\n"
        "• 📷 <b>Понимать фотографии/скриншоты</b>: описание и извлечение важных деталей\n\n"
        "• 🎙️ <b>Распознавать голосовые сообщения</b> и отвечать как текстом, так и голосом\n"
        "• 🗂️ <b>Готовить презентации в PPTX</b> по команде /ppt\n\n"
        "👇 Выбирай, с чего начать:"
    )

    if update.message:
        await update.message.reply_text(text, parse_mode="HTML", reply_markup=main_keyboard())
        try:
            await update.message.reply_text("Быстрые команды доступны на клавиатуре ниже.", reply_markup=QUICK_COMMANDS_KEYBOARD)
        except Exception:
            pass
    else:
        await context.bot.send_message(chat_id=user.id, text=text, parse_mode="HTML", reply_markup=main_keyboard())
        try:
            await context.bot.send_message(chat_id=user.id, text="Быстрые команды доступны на клавиатуре ниже.", reply_markup=QUICK_COMMANDS_KEYBOARD)
        except Exception:
            pass


# =========================
# Профиль
# =========================
async def _render_profile_html(user_id: int) -> str:
    profile = await _ensure_profile(user_id)
    prem = await is_premium(user_id)
    used_today = await get_usage_today(user_id)
    bonus = await get_free_credits(user_id)
    fav_count = len(await list_favorite_prompts(user_id))

    me = await application.bot.get_me()
    deep_link = f"https://t.me/{me.username}?start=ref_{user_id}"
    visual = _user_model_visual.get(user_id, "GPT-4o mini")
    mode_lbl = current_mode_label(user_id)

    if prem:
        left_text = "∞ (Премиум)"
        status = "Премиум"
        # Покажем до какого числа и сколько осталось
        exp_iso = await get_premium_expires(user_id)
        extra = ""
        if exp_iso:
            try:
                exp_dt = datetime.fromisoformat(exp_iso)
            except Exception:
                exp_dt = None
            if exp_dt:
                now_dt = datetime.utcnow()
                if exp_dt.tzinfo:  # если в expires_at есть tz
                    now_dt = datetime.now(exp_dt.tzinfo)
                remaining = exp_dt - now_dt
                days_left = max(0, remaining.days + (1 if remaining.seconds > 0 else 0))
                extra = f"\nПремиум до: <b>{exp_dt.strftime('%d.%m.%Y %H:%M')}</b> (осталось ~<b>{days_left}</b> дн.)"
        status += extra
    else:
        left_day = max(0, DAILY_LIMIT - used_today)
        total_left = left_day + bonus
        left_text = f"{total_left} (дневной лимит {left_day}, бонусов {bonus})"
        status = "Обычный"

    return (
        f"👤 <b>Профиль</b>\n"
        f"ID: <code>{user_id}</code>\n"
        f"Статус: <b>{status}</b>\n"
        f"Осталось заявок: <b>{left_text}</b>\n"
        f"Модель: <b>{visual}</b>\n"
        f"Режим: <b>{mode_lbl}</b>\n\n"
        "🧾 <b>Настройки</b>\n"
        f"• Стиль: <b>{PROFILE_STYLES.get(profile.get('style'), 'Стандарт')}</b>\n"
        f"• Язык: <b>{PROFILE_LANGUAGES.get(profile.get('language'), 'Авто')}</b>\n"
        f"• Формат: <b>{PROFILE_FORMATS.get(profile.get('output_format'), 'Обычный')}</b>\n"
        f"• Тема: <b>{PROFILE_THEMES.get(profile.get('theme'), 'Авто')}</b>\n"
        f"• Избранных шаблонов: <b>{fav_count}</b>\n\n"
        f"🔗 <b>Ваша реферальная ссылка:</b>\n{deep_link}\n\n"
        f"За каждого приглашённого: +{REF_BONUS} заявок."
    )

async def cmd_profile(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    txt = await _render_profile_html(user_id)
    await update.message.reply_text(txt, parse_mode="HTML", reply_markup=main_keyboard())

async def on_profile_btn(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass
    txt = await _render_profile_html(q.from_user.id)
    try:
        await q.message.edit_text(txt, parse_mode="HTML", reply_markup=main_keyboard())
    except Exception:
        await q.message.reply_text(txt, parse_mode="HTML", reply_markup=main_keyboard())

# =========================
# Рефералка
# =========================
async def _render_referral_html(user_id: int) -> str:
    me = await application.bot.get_me()
    deep_link = f"https://t.me/{me.username}?start=ref_{user_id}"
    return (
        "🎁 <b>Реферальная программа</b>\n\n"
        f"Приглашайте друзей по ссылке и получайте <b>+{REF_BONUS}</b> бесплатных заявок за каждого!\n\n"
        f"🔗 Ваша ссылка:\n{deep_link}\n\n"
        "Как это работает:\n"
        "• Человек нажимает по ссылке и жмёт /start\n"
        f"• Вам автоматически начисляется <b>+{REF_BONUS}</b> заявок\n"
        "• Бонусы суммируются и расходуются после дневного лимита\n"
    )

async def cmd_ref(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    txt = await _render_referral_html(user_id)
    await update.message.reply_text(txt, parse_mode="HTML", reply_markup=main_keyboard())

async def on_ref_btn(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass
    txt = await _render_referral_html(q.from_user.id)
    try:
        await q.message.edit_text(txt, parse_mode="HTML", reply_markup=main_keyboard())
    except Exception:
        await q.message.reply_text(txt, parse_mode="HTML", reply_markup=main_keyboard())

# =========================
# Визуальный выбор модели
# =========================
async def on_models_btn(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass
    text = _models_menu_text("short")
    try:
        await q.message.edit_text(text, parse_mode="HTML", reply_markup=models_keyboard_visual())
    except Exception:
        await q.message.reply_text(text, parse_mode="HTML", reply_markup=models_keyboard_visual())

async def on_models_view_toggle(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass
    mode = "short" if q.data == "mvis:short" else "full"
    text = _models_menu_text(mode)
    try:
        await q.message.edit_text(text, parse_mode="HTML", reply_markup=models_keyboard_visual())
    except Exception:
        await q.message.reply_text(text, parse_mode="HTML", reply_markup=models_keyboard_visual())

async def on_model_visual_select(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass
    label = (q.data or "").split("mvis:sel:", 1)[-1].strip() or "GPT-4o mini"

    _user_model_visual[q.from_user.id] = label
    # простая логика: всё, что содержит DeepSeek — на DeepSeek, остальное — OpenAI
    if "DeepSeek" in label:
        _user_model[q.from_user.id] = MODEL_DEEPSEEK
    else:
        _user_model[q.from_user.id] = MODEL_OPENAI

    msg = f"✅ Модель «{label}» установлена.\nМожно писать сообщение!"
    try:
        await q.message.edit_text(msg, reply_markup=main_keyboard())
    except Exception:
        await q.message.reply_text(msg, reply_markup=main_keyboard())

# =========================
# Режимы (ярлыки)
# =========================
async def on_modes_btn(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass
    txt = (
        "Выберите режим ответа:\n"
        "• <b>Стандарт</b> — обычные ответы\n"
        "• <b>Кодинг</b> — больше кода и примеров\n"
        "• <b>SEO</b> — тексты и структура для SEO\n"
        "• <b>Перевод</b> — RU↔EN, аккуратный стиль\n"
        "• <b>Резюме</b> — краткие выжимки\n"
        "• <b>Креатив</b> — идеи, варианты, слоганы"
    )
    try:
        await q.message.edit_text(txt, parse_mode="HTML", reply_markup=modes_keyboard())
    except Exception:
        await q.message.reply_text(txt, parse_mode="HTML", reply_markup=modes_keyboard())

async def on_mode_select(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass
    key = (q.data or "").split("mode:", 1)[-1]
    if key not in TASK_MODES:
        key = "default"
    _user_task_mode[q.from_user.id] = key
    lbl = TASK_MODES[key]["label"]
    try:
        await q.message.edit_text(f"✅ Режим «{lbl}» активирован. Готов работать!", reply_markup=main_keyboard())
    except Exception:
        await q.message.reply_text(f"✅ Режим «{lbl}» активирован. Готов работать!", reply_markup=main_keyboard())

# =========================
# Помощь / FAQ / Оферта
# =========================

def _faq_text() -> str:
    return (
        "<b>FAQ — Частые вопросы</b>\n\n"
        "• <b>Что даёт Премиум?</b>\n"
        "  Безлимитные сообщения, доступ ко всем моделям и генерации изображений.\n\n"
        "• <b>Сколько стоит Премиум и на сколько дней?</b>\n"
        f"  {PRICE_RUB_TEXT} за 30 дней. Оплатить можно через кнопку «Купить подписку».\n\n"
        "• <b>Где посмотреть, когда закончится Премиум?</b>\n"
        "  Откройте «👤 Профиль» — там дата окончания и оставшиеся дни.\n\n"
        "• <b>Как работают лимиты без Премиум?</b>\n"
        f"  {DAILY_LIMIT}/день + реферальные бонусы.\n\n"
        "• <b>Как получить бонусы?</b>\n"
        "  Пригласите друзей по вашей реферальной ссылке из Профиля — за каждого +25 заявок.\n\n"
        "• <b>Могу ли я озвучить ответы бота?</b>\n"
        "  Да, просто нажмите кнопку «🎧 Озвучить» под любым сообщением.\n\n"
        "• <b>Можно ли отправлять документы?</b>\n"
        "  Да, бот поддерживает .txt, .md, .csv и .pdf — он сделает краткое резюме и выделит ключевые пункты.\n\n"
        "• <b>Можно ли анализировать фото?</b>\n"
        "  Да, просто отправьте изображение или скриншот — бот опишет, что на нём, и выделит детали.\n\n"
        "• <b>Возвраты и вопросы по оплатам</b>\n"
        f"  Пишите на <a href='mailto:{SUPPORT_EMAIL}'>{SUPPORT_EMAIL}</a> — поможем. "
        "Возврат средств возможен в случаях, предусмотренных законодательством РФ и условиями нашей оферты.\n\n"
        "• <b>Конфиденциальность</b>\n"
        "  Мы обрабатываем персональные данные по ФЗ-152 и используем их только для оказания услуг.\n"
    )

def _support_text() -> str:
    return (
        "<b>Техподдержка</b>\n\n"
        f"• Email: <a href='mailto:{SUPPORT_EMAIL}'>{SUPPORT_EMAIL}</a>\n"
        f"• Время ответа: в рабочие часы {SUPPORT_WORK_HOURS}\n\n"
        "В письме укажите: ID в Telegram (из Профиля), кратко суть вопроса, скрин/ошибку и время события."
    )

async def on_help_btn(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass
    kb = InlineKeyboardMarkup([
        [InlineKeyboardButton("📚 FAQ", callback_data="help:faq"),
         InlineKeyboardButton("🛟 Техподдержка", callback_data="help:support")],
        [InlineKeyboardButton("📄 Публичная оферта", url=PUBLIC_OFFER_URL)],
        [InlineKeyboardButton("✉️ Написать на email", url=f"mailto:{SUPPORT_EMAIL}")],
        [InlineKeyboardButton("⬅️ Назад", callback_data="home")]
    ])
    await q.message.edit_text("Раздел помощи. Выберите:", reply_markup=kb)

async def on_help_faq(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass
    kb = InlineKeyboardMarkup([
        [InlineKeyboardButton("🛟 Техподдержка", callback_data="help:support")],
        [InlineKeyboardButton("📄 Публичная оферта", url=PUBLIC_OFFER_URL)],
        [InlineKeyboardButton("⬅️ Назад", callback_data="home")]
    ])
    await q.message.edit_text(_faq_text(), parse_mode="HTML", reply_markup=kb)

async def on_help_support(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass
    kb = InlineKeyboardMarkup([
        [InlineKeyboardButton("✉️ Написать на email", url=f"mailto:{SUPPORT_EMAIL}")],
        [InlineKeyboardButton("📚 FAQ", callback_data="help:faq")],
        [InlineKeyboardButton("📄 Публичная оферта", url=PUBLIC_OFFER_URL)],
        [InlineKeyboardButton("⬅️ Назад", callback_data="home")]
    ])
    await q.message.edit_text(_support_text(), parse_mode="HTML", reply_markup=kb)

async def on_help_how(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Открывает тот же текст, что и /help, но по кнопке."""
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass

    txt = (
        "<b>Как пользоваться ботом</b>\n\n"
        "• Напишите любое сообщение — я отвечу.\n"
        "• 🎧 Озвучивай ответы: нажми кнопку «Озвучить» под сообщением.\n"
        "• 📄 Отправляй документы (.txt, .md, .csv, .pdf) — сделаю краткое резюме.\n"
        "• 📷 Присылай фото или скриншоты — опишу, что на них.\n"
        "• Нужна картинка? Команда /img.\n"
        "• 🗂️ Генерируй презентации — /ppt <тема>.\n"
        "• Выбор модели — /models.\n"
        "• Переключить режим — /mode.\n"
        "• Профиль и рефералка — /profile, /ref.\n"
       f"• Премиум ({PRICE_RUB_TEXT} / 30 дней) — /buy.\n\n"
        "Кнопки ниже — быстрый доступ:"
    )

    # мини-меню помощи с быстрыми ссылками
    kb = InlineKeyboardMarkup([
        [InlineKeyboardButton("📚 FAQ",          callback_data="help:faq"),
         InlineKeyboardButton("🛟 Техподдержка", callback_data="help:support")],
        [InlineKeyboardButton("📄 Публичная оферта", url=PUBLIC_OFFER_URL)],
        [InlineKeyboardButton("⬅️ Назад", callback_data="home")]
    ])
    await q.message.edit_text(txt, parse_mode="HTML", reply_markup=kb)

# Команды-псевдонимы
async def cmd_support(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text(_support_text(), parse_mode="HTML", reply_markup=InlineKeyboardMarkup([
        [InlineKeyboardButton("✉️ Написать на email", url=f"mailto:{SUPPORT_EMAIL}")],
        [InlineKeyboardButton("📚 FAQ", callback_data="help:faq")],
        [InlineKeyboardButton("📄 Публичная оферта", url=PUBLIC_OFFER_URL)],
        [InlineKeyboardButton("⬅️ Назад", callback_data="home")]
    ]))

async def cmd_faq(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text(_faq_text(), parse_mode="HTML", reply_markup=InlineKeyboardMarkup([
        [InlineKeyboardButton("🛟 Техподдержка", callback_data="help:support")],
        [InlineKeyboardButton("📄 Публичная оферта", url=PUBLIC_OFFER_URL)],
        [InlineKeyboardButton("⬅️ Назад", callback_data="home")]
    ]))

async def on_ppt_btn(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass
    text = (
        "🗂️ <b>Генерация презентаций</b>\n\n"
        "Отправьте команду вида:\n"
        "<code>/ppt тема презентации</code>\n\n"
        "Например: <code>/ppt План вывода нового продукта</code>.\n"
        "Доступно для пользователей с Премиум-подпиской."
    )
    try:
        await q.message.edit_text(text, parse_mode="HTML")
    except Exception:
        await q.message.reply_text(text, parse_mode="HTML")

# =========================
# Оплата (CryptoPay)
# =========================
async def _create_crypto_invoice_link(user_id: int) -> str:
    if not CRYPTOPAY_KEY:
        raise RuntimeError("Оплата не подключена (нет CRYPTOPAY_KEY).")

    payload = str(user_id)
    headers = {"Crypto-Pay-API-Token": CRYPTOPAY_KEY}
    data = {
        "asset": "USDT",
        "amount": PRICE_USDT,
        "description": f"Подписка на 30 дней ({PRICE_RUB_TEXT})",
        "payload": payload,
    }

    try:
        async with httpx.AsyncClient(timeout=15.0) as client:
            response = await client.post(
                "https://pay.crypt.bot/api/createInvoice",
                json=data,
                headers=headers,
            )
        response.raise_for_status()
    except httpx.HTTPStatusError as exc:
        detail = exc.response.text[:400] if exc.response is not None else str(exc)
        raise RuntimeError(f"CryptoPay вернул ошибку: {detail}") from exc
    except httpx.RequestError as exc:
        raise RuntimeError(f"CryptoPay недоступен: {exc}") from exc

    try:
        payload_json = response.json()
    except ValueError as exc:
        raise RuntimeError("CryptoPay вернул некорректный JSON.") from exc

    result = payload_json.get("result") or {}
    url = result.get("pay_url")
    if not url:
        raise RuntimeError(f"Не удалось получить ссылку оплаты: {payload_json}")
    return url


async def on_buy_btn(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    try:
        await q.answer()
    except Exception:
        pass

    if not CRYPTOPAY_KEY:
        await q.message.reply_text("Оплата не подключена (нет CRYPTOPAY_KEY).")
        return

    try:
        url = await _create_crypto_invoice_link(q.from_user.id)
    except Exception as e:
        await q.message.reply_text(f"Не удалось создать счёт: {e}")
        return

    text = (
    f"Оплата подписки на 30 дней: <b>{PRICE_RUB_TEXT}</b>\n\n"
    "<b>Премиум даёт</b>:\n"
    "• Безлимитные сообщения (без очередей)\n"
    "• Доступ ко <b>всем</b> моделям\n"
    "• <b>Генерацию изображений</b> (Replicate · Flux-1 Schnell)\n"
    "• Приоритетную обработку\n\n"
    f"Ссылка на оплату:\n{url}"
    )
    await q.message.reply_text(text, parse_mode="HTML")

# =========================
# Команды /buy /models /mode /help
# =========================
async def cmd_buy(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Команда /buy — та же логика, что и по кнопке."""
    if not CRYPTOPAY_KEY:
        await update.message.reply_text("Оплата не подключена (нет CRYPTOPAY_KEY).")
        return

    try:
        url = await _create_crypto_invoice_link(update.effective_user.id)
    except Exception as e:
        await update.message.reply_text(f"Не удалось создать счёт: {e}")
        return

    text = (
    f"Оплата подписки на 30 дней: <b>{PRICE_RUB_TEXT}</b>\n\n"
    "<b>Премиум даёт</b>:\n"
    "• Безлимитные сообщения (без очередей)\n"
    "• Доступ ко <b>всем</b> моделям\n"
    "• <b>Генерацию изображений</b> (Replicate · Flux-1 Schnell)\n"
    "• Приоритетную обработку\n\n"
    f"Ссылка на оплату:\n{url}"
    )
    await update.message.reply_text(text, parse_mode="HTML")

async def cmd_models(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Команда /models — показать меню выбора модели."""
    text = _models_menu_text("short")
    await update.message.reply_text(text, parse_mode="HTML", reply_markup=models_keyboard_visual())

async def cmd_mode(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Команда /mode — показать меню режимов."""
    txt = (
        "Выберите режим ответа:\n"
        "• <b>Стандарт</b> — обычные ответы\n"
        "• <b>Кодинг</b> — больше кода и примеров\n"
        "• <b>SEO</b> — тексты и структура для SEO\n"
        "• <b>Перевод</b> — RU↔EN, аккуратный стиль\n"
        "• <b>Резюме</b> — краткие выжимки\n"
        "• <b>Креатив</b> — идеи, варианты, слоганы"
    )
    await update.message.reply_text(txt, parse_mode="HTML", reply_markup=modes_keyboard())

async def cmd_help(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Команда /help — краткая справка с основными действиями."""
    txt = (
        "<b>Как пользоваться ботом</b>\n\n"
        "• Напишите любое сообщение — я отвечу.\n"
        "• 🎧 Озвучивай ответы: нажми кнопку «Озвучить» под сообщением.\n"
        "• 📄 Отправляй документы (.txt, .md, .csv, .pdf) — я сделаю краткое резюме.\n"
        "• 📷 Присылай фото или скриншоты — расскажу, что на них.\n"
        "• Нужна картинка? Команда /img.\n"
        "• 🗂️ Генерируй презентации — /ppt <тема>.\n"
        "• Выбор модели — /models.\n"
        "• Переключить режим — /mode.\n"
        "• Профиль и рефералка — /profile, /ref.\n"
       f"• Премиум ({PRICE_RUB_TEXT} / 30 дней) — /buy.\n\n"
        "Кнопки ниже — быстрый доступ:"
    )
    await update.message.reply_text(txt, parse_mode="HTML", reply_markup=main_keyboard())

async def cmd_ppt(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Команда /ppt — сгенерировать PPTX по теме."""
    user_id = update.effective_user.id
    topic = " ".join(context.args).strip() if context.args else ""
    if not topic:
        await update.message.reply_text(
            "Использование: /ppt <тема/задача>.\n"
            "Например: /ppt маркетинговая стратегия для нового продукта."
        )
        return

    if not await is_premium(user_id):
        await update.message.reply_text(
            "Генерация презентаций доступна только в Премиум.",
            reply_markup=InlineKeyboardMarkup(
                [[InlineKeyboardButton("💳 Купить подписку", callback_data="buy")]]
            ),
        )
        return

    status = await update.message.reply_text("🧩 Составляю структуру презентации…")
    ppt_path: Path | None = None
    hero_image: Path | None = None
    try:
        slides = _generate_presentation_structure(user_id, topic)
        if not slides:
            raise RuntimeError("Не удалось собрать структуру презентации.")

        tmpdir = Path(tempfile.gettempdir())
        ppt_path = tmpdir / f"presentation_{user_id}_{int(time.time())}.pptx"
        palette = _choose_color_palette(user_id, topic)
        hero_image = await _generate_presentation_image(topic)
        _build_presentation_file(slides, ppt_path, topic, palette, hero_image)

        try:
            await status.edit_text("📤 Отправляю файл…")
        except Exception:
            pass

        safe_name = re.sub(r"[^A-Za-z0-9]+", "_", topic)[:40] or "presentation"
        with open(ppt_path, "rb") as doc:
            await update.message.reply_document(
                document=doc,
                filename=f"{safe_name}.pptx",
                caption="Презентация готова ✅",
            )
    except Exception as e:
        try:
            await status.edit_text(f"Не удалось создать презентацию: {e}")
        except Exception:
            await update.message.reply_text(f"Не удалось создать презентацию: {e}")
        return
    finally:
        try:
            await status.delete()
        except Exception:
            pass
        if ppt_path:
            try:
                ppt_path.unlink(missing_ok=True)
            except Exception:
                pass
        if hero_image:
            try:
                hero_image.unlink(missing_ok=True)
            except Exception:
                pass

# =========================
# Сообщения пользователей
# =========================
async def _handle_text_request(update: Update, context: ContextTypes.DEFAULT_TYPE, text: str):
    user_id = update.effective_user.id
    text = text or ""

    if not text.strip():
        await update.message.reply_text("Сообщение пустое. Пожалуйста, отправьте текст.")
        return

    await _ensure_profile(user_id)

    # Переименование чата — если ждём от пользователя новое имя
    if _pending_chat_rename.get(user_id):
        cid = _pending_chat_rename[user_id]
        new_title = text.strip()[:80]
        if not new_title:
            await update.message.reply_text(
                "Название пустое. Отправьте текст от 1 до 80 символов или нажмите «Отмена» в меню."
            )
            return
        ok = await rename_chat(user_id, cid, new_title)
        _pending_chat_rename.pop(user_id, None)
        if ok:
            await update.message.reply_text("Готово: чат переименован ✅")
        else:
            await update.message.reply_text("Не удалось переименовать чат.")
        return

    # Если ждём промпт для изображения — перехватываем
    if _awaiting_img_prompt.get(user_id):
        _awaiting_img_prompt[user_id] = False
        if not await is_premium(user_id):
            await update.message.reply_text(
                "Генерация изображений только для Премиум.",
                reply_markup=InlineKeyboardMarkup(
                    [[InlineKeyboardButton("Купить подписку", callback_data="buy")]]
                ),
            )
            return
        await update.message.reply_text("Генерирую изображение…")
        await generate_image_and_send(user_id, update.effective_chat.id, text, context.bot)
        return

    _last_user_prompt[user_id] = text

    # лимиты как раньше
    if not await is_premium(user_id):
        if await can_send_message(user_id, limit=DAILY_LIMIT):
            pass
        elif await consume_free_credit(user_id):
            pass
        else:
            await update.message.reply_text(
                "🚫 Лимит исчерпан.\n"
                f"— Дневной лимит: {DAILY_LIMIT}/день\n"
                f"— Реферальные бонусы: получите +{REF_BONUS} заявок за каждого приглашённого!\n\n"
                "Купите подписку «💳 Купить подписку» для безлимита."
            )
            return

    # выбор по диалоговому режиму
    mode = await get_chat_mode(user_id)

    # ➊ Поставим временное сообщение
    spinner = await update.message.reply_text("🤖 Генерация ответа…")

    try:
        if mode == DIALOG_ROOMS:
            # нужен активный чат, если нет — создадим
            cid = await get_active_chat(user_id)
            if cid is None:
                cid = await create_chat(user_id, "Чат 1")
                await set_active_chat(user_id, cid)

            # загрузим историю (последние 20 сообщений) + добавим текущий запрос
            history = await get_chat_history(cid, limit=20)
            reply = ask_llm_context(user_id, history, text)

            # сохраним и вопрос, и ответ в историю
            await add_chat_message(cid, "user", text)
            await add_chat_message(cid, "assistant", reply)
        else:
            # быстрый режим
            reply = ask_llm(user_id, text)
    finally:
        # ➋ Удаляем «спиннер» в любом исходе
        try:
            await context.bot.delete_message(
                chat_id=update.effective_chat.id,
                message_id=spinner.message_id,
            )
        except Exception:
            pass

    # ➌ Отправляем ответ (кнопки как обсуждали)
    _last_answer[user_id] = reply
    parts = _split_for_telegram(reply)
    buttons: list[list[InlineKeyboardButton]] = []
    if len(parts) == 1:
        _last_answer[user_id] = parts[0]
        buttons.append([InlineKeyboardButton("🎧 Озвучить", callback_data="tts")])
    else:
        _long_reply_queue[user_id] = parts[1:]
        _last_answer[user_id] = parts[0]
        buttons.append([
            InlineKeyboardButton("Показать ещё ▶️", callback_data="more"),
            InlineKeyboardButton("🎧 Озвучить", callback_data="tts"),
        ])
    buttons.append([
        InlineKeyboardButton("⭐ Шаблон", callback_data="fav:add"),
        InlineKeyboardButton("🔁 Перевести", callback_data="quick:translate"),
        InlineKeyboardButton("🧾 Сжать", callback_data="quick:summary"),
    ])
    await update.message.reply_text(parts[0], reply_markup=InlineKeyboardMarkup(buttons))
    return


async def on_message(update: Update, context: ContextTypes.DEFAULT_TYPE):
    text = update.message.text or ""
    await _handle_text_request(update, context, text)

async def on_voice(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    status = await update.message.reply_text("🎙️ Распознаю голос…")

    tmpdir = Path(tempfile.gettempdir())
    audio_path = tmpdir / f"voice_{user_id}_{int(time.time())}.ogg"

    try:
        data = await _download_telegram_file(context.bot, update.message.voice.file_id)
        with open(audio_path, "wb") as f:
            f.write(data)
        transcript = await asyncio.to_thread(_transcribe_audio_file_sync, audio_path)
        transcript = (transcript or "").strip()
    except Exception as e:
        await status.edit_text(f"Не удалось распознать голос: {e}")
        return
    finally:
        try:
            audio_path.unlink(missing_ok=True)
        except Exception:
            pass

    if not transcript:
        await status.edit_text("Не удалось распознать голосовое сообщение.")
        return

    try:
        await status.edit_text(f"📝 Распознано: {transcript}")
    except Exception:
        pass

    await _handle_text_request(update, context, transcript)

# =========================
# Обработчик фото
# =========================
async def on_photo(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
        # --- КД на отправку фото (для всех, включая премиум) ---
    now = time.time()
    until = _photo_cd_until.get(user_id, 0)
    if until > now:
        left = int(until - now)
        await update.message.reply_text(f"⏳ Подождите {left} сек. перед следующей фотографией.")
        return
    # ставим новый КД сразу, чтобы не спамили даже при ошибке
    _photo_cd_until[user_id] = now + PHOTO_COOLDOWN_SEC
    # лимиты как в on_message
    if not await is_premium(user_id):
        if await can_send_message(user_id, limit=DAILY_LIMIT):
            pass
        elif await consume_free_credit(user_id):
            pass
        else:
            await update.message.reply_text(
                "🚫 Лимит исчерпан.\n"
                f"— Дневной лимит: {DAILY_LIMIT}/день\n"
                f"— Реферальные бонусы: получите +{REF_BONUS} заявок за каждого приглашённого!\n\n"
                "Купите подписку «💳 Купить подписку» для безлимита."
            )
            return

    spinner = await update.message.reply_text("🤖 Генерация ответа…")
    try:
        # берём самую большую превьюху
        photo = update.message.photo[-1]
        data = await _download_telegram_file(context.bot, photo.file_id)
        img64 = _img_b64(data)
        # если у сообщения есть подпись — используем как hint
        hint = update.message.caption or ""
        reply = _analyze_image_with_llm(user_id, hint, img64)
    except Exception as e:
        await update.message.reply_text(f"Не удалось проанализировать изображение: {e}")
        return
    finally:
        try:
            await context.bot.delete_message(
                chat_id=update.effective_chat.id,
                message_id=spinner.message_id,
            )
        except Exception:
            pass

    _last_answer[user_id] = reply
    chunks = _split_for_telegram(reply)
    if len(chunks) > 1:
        _long_reply_queue[user_id] = chunks[1:]
        kb = InlineKeyboardMarkup(
            [
                [
                    InlineKeyboardButton("Показать ещё ▶️", callback_data="more"),
                    InlineKeyboardButton("🎧 Озвучить", callback_data="tts"),
                ]
            ]
        )
        await update.message.reply_text(chunks[0], reply_markup=kb)
    else:
        kb = InlineKeyboardMarkup([[InlineKeyboardButton("🎧 Озвучить", callback_data="tts")]])
        await update.message.reply_text(chunks[0], reply_markup=kb)

# =========================
# Обработчик документов (.txt/.md/.csv/.pdf)
# =========================
async def on_document(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id

    # лимиты
    if not await is_premium(user_id):
        if await can_send_message(user_id, limit=DAILY_LIMIT):
            pass
        elif await consume_free_credit(user_id):
            pass
        else:
            await update.message.reply_text(
                "🚫 Лимит исчерпан.\n"
                f"— Дневной лимит: {DAILY_LIMIT}/день\n"
                f"— Реферальные бонусы: получите +{REF_BONUS} заявок!\n\n"
                "Купите подписку «💳 Купить подписку» для безлимита."
            )
            return

    doc = update.message.document
    title = doc.file_name or "документ"
    spinner = await update.message.reply_text("🤖 Генерация ответа…")
    try:
        data = await _download_telegram_file(context.bot, doc.file_id)
        text_content = ""
        lower = (title or "").lower()

        if lower.endswith((".txt", ".md", ".csv")):
            # простые текстовые — читаем как utf-8
            text_content = data.decode("utf-8", errors="replace")
        elif lower.endswith(".pdf"):
            import io
            reader = PdfReader(io.BytesIO(data))
            pages = min(10, len(reader.pages))  # не больше 10 страниц
            chunks = []
            for i in range(pages):
                try:
                    chunks.append(reader.pages[i].extract_text() or "")
                except Exception:
                    pass
            text_content = "\n\n".join(chunks).strip()
            if not text_content:
                text_content = "[Не удалось извлечь текст из PDF. Попробуйте отправить как изображение/скриншот.]"
        else:
            await update.message.reply_text("Поддерживаю пока .txt, .md, .csv и .pdf. Попробуйте один из этих форматов.")
            return

        reply = _summarize_text_with_llm(user_id, title, text_content)
        try:
            await context.bot.delete_message(chat_id=update.effective_chat.id,
                                            message_id=spinner.message_id)
        except Exception:
            pass
        _last_answer[user_id] = reply
        chunks = _split_for_telegram(reply)
        if len(chunks) > 1:
            _long_reply_queue[user_id] = chunks[1:]
            kb = InlineKeyboardMarkup([
                [InlineKeyboardButton("Показать ещё ▶️", callback_data="more"),
                 InlineKeyboardButton("🎧 Озвучить", callback_data="tts")]
            ])
            await update.message.reply_text(chunks[0], reply_markup=kb)
        else:
            kb = InlineKeyboardMarkup([[InlineKeyboardButton("🎧 Озвучить", callback_data="tts")]])
            await update.message.reply_text(chunks[0], reply_markup=kb)

    except Exception as e:
        await update.message.reply_text(f"Не удалось обработать документ: {e}")

# =========================
# Админ-команды
# =========================
async def cmd_admin(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ADMIN_ID:
        await update.message.reply_text("⛔ Нет доступа.")
        return
    paid_today = await count_paid_users_today()
    paid_total = await count_paid_users_total()
    await update.message.reply_text(
        "📊 Админ-панель\n"
        f"Покупок сегодня: {paid_today}\n"
        f"Всего активных премиумов: {paid_total}\n\n"
        "Команды:\n"
        "/add_premium <user_id> <days>\n"
        "/remove_premium <user_id>\n"
        "/broadcast <text>"
    )

async def cmd_add_premium(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ADMIN_ID:
        return
    if len(context.args) < 2:
        await update.message.reply_text("Формат: /add_premium <user_id> <days>")
        return
    try:
        uid = int(context.args[0])
        days = int(context.args[1])
        expires_at = (datetime.now() + timedelta(days=days)).isoformat()
        await set_premium(uid, expires_at)
        await update.message.reply_text(f"✅ Премиум выдан {uid} на {days} дн.")
        try:
            await application.bot.send_message(uid, f"🎉 Вам выдали премиум на {days} дней!")
        except Exception:
            pass
    except Exception as e:
        await update.message.reply_text(f"Ошибка: {e}")

async def cmd_remove_premium(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ADMIN_ID:
        return
    try:
        if not context.args:
            await update.message.reply_text("Формат: /remove_pремиum <user_id>")
            return
        uid = int(context.args[0])
        await revoke_premium(uid)
        await update.message.reply_text(f"❎ Премиум снят у {uid}.")
        try:
            await application.bot.send_message(uid, "⚠️ Ваш премиум был отключён.")
        except Exception:
            pass
    except Exception as e:
        await update.message.reply_text(f"Ошибка: {e}")

async def cmd_broadcast(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ADMIN_ID:
        return
    if not context.args:
        await update.message.reply_text("Формат: /broadcast <text>")
        return
    text = " ".join(context.args)
    await update.message.reply_text(f"Ок, отправлю: {text}\n(реальную рассылку можно дописать в db.py)")

# =========================
# Webhooks
# =========================
@app.post("/tg")
async def telegram_webhook(request: Request):
    global application
    if application is None:
        return {"ok": False, "error": "bot not initialized"}
    data = await request.json()
    update = Update.de_json(data, application.bot)
    uid = update.update_id
    if uid in _recent_set:
        return {"ok": True}  # уже обрабатывали
    _recent_updates.append(uid)
    _recent_set.add(uid)
    if len(_recent_updates) == _recent_updates.maxlen:
    # чистим множество, когда очередь заполнена
        try:
            while len(_recent_set) > _recent_updates.maxlen:
                _recent_set.remove(_recent_updates.popleft())
        except Exception:
            _recent_set.clear()
            _recent_set.update(_recent_updates)
    await application.process_update(update)
    return {"ok": True}

@app.post("/cryptopay-webhook")
async def cryptopay_webhook(request: Request):
    """Обработчик вебхуков Crypto Pay (update_type=invoice_paid)."""
    global application
    if not CRYPTOPAY_KEY:
        return {"ok": False, "error": "cryptopay disabled"}

    try:
        raw_body = await request.body()
    except Exception:
        return {"ok": False, "error": "bad body"}

    signature = (
        request.headers.get("Crypto-Pay-Signature")
        or request.headers.get("X-Crypto-Pay-Signature")
        or request.headers.get("X-CryptoPay-Signature")
    )
    if not signature:
        logger.warning("CryptoPay webhook: missing signature header")
        return {"ok": False, "error": "signature missing"}

    expected_sig = hmac.new(
        CRYPTOPAY_KEY.encode("utf-8"),
        raw_body,
        hashlib.sha256,
    ).hexdigest()
    if not hmac.compare_digest(signature.strip().lower(), expected_sig):
        logger.warning("CryptoPay webhook: invalid signature")
        return {"ok": False, "error": "invalid signature"}

    try:
        data = json.loads(raw_body.decode("utf-8"))
    except Exception:
        return {"ok": False, "error": "bad json"}

    try:
        logger.info("CryptoPay webhook: %s", data)
    except Exception:
        pass

    user_id = None
    paid = False

    # Новый формат
    update_type = data.get("update_type")
    inv_new = data.get("payload") or {}
    if update_type == "invoice_paid" and isinstance(inv_new, dict):
        raw_uid = inv_new.get("payload")
        status_new = inv_new.get("status")
        if raw_uid is not None and (status_new is None or status_new == "paid"):
            try:
                user_id = int(str(raw_uid))
                paid = True
            except Exception:
                user_id = None

    # Совместимость со старым форматом
    if not paid:
        invoice = data.get("invoice") or {}
        status = invoice.get("status")
        raw_uid = invoice.get("payload")
        if status == "paid" and raw_uid is not None:
            try:
                user_id = int(str(raw_uid))
                paid = True
            except Exception:
                user_id = None

    if paid and user_id:
        expires_dt = datetime.now() + timedelta(days=30)
        await set_premium(user_id, expires_dt.isoformat())
        try:
            text = (
                "✅ <b>Оплата получена</b>!\n"
                f"Премиум активирован до <b>{expires_dt.strftime('%d.%m.%Y')}</b>.\n\n"
                "Что дальше?\n"
                "• Откройте профиль — проверить статус и реф. ссылку\n"
                "• Выберите модель — переключиться на нужный режим\n"
                "• Или просто напишите сообщение 🙂"
            )
            await application.bot.send_message(
                chat_id=user_id,
                text=text,
                parse_mode="HTML",
                reply_markup=main_keyboard()
            )
        except Exception:
            try:
                await application.bot.send_message(
                    chat_id=user_id,
                    text="✅ Оплата получена! Премиум активирован на 30 дней."
                )
            except Exception:
                pass

    return {"ok": True}

@app.get("/share/{token}")
async def share_chat(token: str):
    data = await get_chat_share(token)
    if not data:
        return PlainTextResponse("Link expired or invalid.", status_code=404)
    user_id, chat_id = data
    history = await get_chat_history_all(chat_id)
    title, _ = await _get_chat_meta(user_id, chat_id)
    body = _chat_history_to_html(title, history) if history else "<p>Чат пуст.</p>"
    html_page = f"""
    <html>
    <head>
        <meta charset="utf-8" />
        <title>{html.escape(title)} · NeuroBot</title>
        <style>
            body {{ font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif; background:#f4f5fb; color:#1f2333; margin:0; padding:40px; }}
            .card {{ max-width: 920px; margin:0 auto; background:white; border-radius:18px; padding:32px; box-shadow:0 14px 35px rgba(31,35,51,0.08); }}
            h1 {{ margin-top:0; font-size:32px; }}
            .entry {{ border-top:1px solid #E5E8F0; padding:18px 0; }}
            .entry:first-of-type {{ border-top:none; }}
            .meta {{ font-size:13px; color:#63708f; margin-bottom:6px; text-transform:uppercase; letter-spacing:0.04em; }}
            .content {{ font-size:16px; line-height:1.6; white-space:pre-wrap; }}
        </style>
    </head>
    <body>
        <div class="card">
            {body}
        </div>
    </body>
    </html>
    """
    return HTMLResponse(html_page)

@app.get("/health")
async def health():
    return {"status": "ok", "time": datetime.utcnow().isoformat()}

# =========================
# Keep-alive (40s) + авто-починка вебхука
# =========================
def _keepalive_loop():
    if not _public_url:
        return
    url = f"{_public_url.rstrip('/')}/health"
    session = requests.Session()
    while not _keepalive_stop.wait(40):
        try:
            session.get(url, timeout=8)
        except Exception:
            pass

async def _webhook_guard_loop():
    """Раз в 10 минут проверяем webhook и чиним, если он слетел."""
    await asyncio.sleep(8)
    while True:
        try:
            bot = application.bot
            _ = await bot.get_me()
            info = await bot.get_webhook_info()
            needed = f"{_public_url.rstrip('/')}/tg"
            if info.url != needed:
                try:
                    await bot.set_webhook(needed, max_connections=40, drop_pending_updates=False)
                    logger.info("🔧 Webhook восстановлен: %s", needed)
                except Exception as e:
                    logger.warning("Webhook repair failed: %s", e)
        except Exception as e:
            logger.warning("webhook guard error: %s", e)
        await asyncio.sleep(600)  # 10 минут

async def _premium_expiry_notifier_loop():
    """Раз в 15 минут ищем истёкшие премиумы и шлём 1 уведомление."""
    await asyncio.sleep(10)
    while True:
        try:
            now_iso = datetime.utcnow().isoformat()
            user_ids = await list_expired_unnotified(now_iso)
            for uid in user_ids:
                # отправим уведомление
                try:
                    await application.bot.send_message(
                        chat_id=uid,
                        text=(
                            "⛔️ Ваш премиум закончился.\n\n"
                            "Продлите подписку, чтобы сохранить безлимит, доступ ко всем моделям "
                            "и генерацию изображений."
                        ),
                        reply_markup=InlineKeyboardMarkup(
                            [[InlineKeyboardButton("💳 Купить подписку", callback_data="buy")]]
                        )
                    )
                except Exception:
                    pass
                # пометить, что уведомили
                try:
                    await mark_expired_notified(uid, now_iso)
                except Exception:
                    pass
            try:
                await cleanup_chat_shares(now_iso)
            except Exception:
                pass
        except Exception as e:
            logger.warning("premium notifier error: %s", e)
        await asyncio.sleep(900)  # 15 минут

# =========================
# Глобальный error-handler PTB (чтобы не падал на 400)
# =========================
async def on_error(update: object, context: ContextTypes.DEFAULT_TYPE):
    logger.warning("PTB error: %s", getattr(context, "error", None))

# =========================
# Инициализация
# =========================
def build_application() -> Application:
    app_ = ApplicationBuilder().token(BOT_TOKEN).build()

    # команды
    app_.add_handler(CommandHandler("start",   cmd_start))
    app_.add_handler(CommandHandler("profile", cmd_profile))
    app_.add_handler(CommandHandler("ref",     cmd_ref))
    app_.add_handler(CommandHandler("admin",   cmd_admin))
    app_.add_handler(CommandHandler("add_premium",    cmd_add_premium))
    app_.add_handler(CommandHandler("remove_premium", cmd_remove_premium))
    app_.add_handler(CommandHandler("broadcast",      cmd_broadcast))
    app_.add_handler(CommandHandler("buy",    cmd_buy))
    app_.add_handler(CommandHandler("models", cmd_models))
    app_.add_handler(CommandHandler("mode",   cmd_mode))
    app_.add_handler(CommandHandler("settings", cmd_settings))
    app_.add_handler(CommandHandler("help",   cmd_help))
    app_.add_handler(CommandHandler("favorites", cmd_favorites))
    app_.add_handler(CommandHandler("ppt",    cmd_ppt))
    app_.add_handler(CommandHandler("support", cmd_support))
    app_.add_handler(CommandHandler("faq",     cmd_faq))

    # кнопка/команда генерации изображений
    app_.add_handler(CallbackQueryHandler(on_img_btn, pattern=r"^img$"))
    app_.add_handler(CommandHandler("img", cmd_img))
    app_.add_handler(CallbackQueryHandler(on_tts_btn, pattern=r"^tts$"))
    app_.add_handler(CallbackQueryHandler(on_more_btn, pattern=r"^more$"))
    app_.add_handler(CallbackQueryHandler(on_quick_translate, pattern=r"^quick:translate$"))
    app_.add_handler(CallbackQueryHandler(on_quick_summary, pattern=r"^quick:summary$"))
    app_.add_handler(CallbackQueryHandler(on_fav_add, pattern=r"^fav:add$"))
    app_.add_handler(CallbackQueryHandler(on_favorites_btn, pattern=r"^fav:list$"))
    app_.add_handler(CallbackQueryHandler(on_fav_run, pattern=r"^fav:run:\d+$"))
    app_.add_handler(CallbackQueryHandler(on_fav_delete, pattern=r"^fav:del:\d+$"))

    # кнопки
    app_.add_handler(CallbackQueryHandler(on_buy_btn,      pattern=r"^buy$"))
    app_.add_handler(CallbackQueryHandler(on_profile_btn,  pattern=r"^profile$"))
    app_.add_handler(CallbackQueryHandler(on_settings_btn, pattern=r"^settings$"))
    app_.add_handler(CallbackQueryHandler(on_settings_change, pattern=r"^settings:(style|language|format|theme):.+$"))
    app_.add_handler(CallbackQueryHandler(on_ref_btn,      pattern=r"^ref$"))
    app_.add_handler(CallbackQueryHandler(on_models_btn,   pattern=r"^models$"))
    app_.add_handler(CallbackQueryHandler(on_models_view_toggle, pattern=r"^mvis:(short|full)$"))
    app_.add_handler(CallbackQueryHandler(on_model_visual_select, pattern=r"^mvis:sel:.+$"))
    app_.add_handler(CallbackQueryHandler(on_modes_btn,    pattern=r"^modes$"))
    app_.add_handler(CallbackQueryHandler(on_mode_select,  pattern=r"^mode:(default|coding|seo|translate|summarize|creative)$"))
    app_.add_handler(CallbackQueryHandler(on_ppt_btn,      pattern=r"^ppt$"))
    app_.add_handler(CallbackQueryHandler(
        lambda u, c: u.callback_query.message.edit_text("Главное меню:", reply_markup=main_keyboard()),
        pattern=r"^home$"
    ))
   
    # помощь / faq
    app_.add_handler(CallbackQueryHandler(on_help_btn,     pattern=r"^help$"))
    app_.add_handler(CallbackQueryHandler(on_help_how,     pattern=r"^help:how$"))   # ← NEW
    app_.add_handler(CallbackQueryHandler(on_help_faq,     pattern=r"^help:faq$"))
    app_.add_handler(CallbackQueryHandler(on_help_support, pattern=r"^help:support$"))

    # сообщения
    app_.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, on_message))
    app_.add_handler(MessageHandler(filters.VOICE & ~filters.COMMAND, on_voice))
    # вложения
    app_.add_handler(MessageHandler(filters.PHOTO & ~filters.COMMAND, on_photo))
    app_.add_handler(MessageHandler(filters.Document.ALL & ~filters.COMMAND, on_document))
    
    # диалоговые режимы и чаты
    app_.add_handler(CallbackQueryHandler(on_dialog_btn,    pattern=r"^dialog$"))
    app_.add_handler(CallbackQueryHandler(on_dialog_select, pattern=r"^dialog:(simple|rooms)$"))
    app_.add_handler(CallbackQueryHandler(on_chats_btn,     pattern=r"^chats$"))
    app_.add_handler(CallbackQueryHandler(on_chat_new,      pattern=r"^chat:new$"))
    app_.add_handler(CallbackQueryHandler(on_chat_open,     pattern=r"^chat:open:\d+$"))
    app_.add_handler(CallbackQueryHandler(on_chat_rename_ask,   pattern=r"^chat:rename:\d+$"))
    app_.add_handler(CallbackQueryHandler(on_chat_delete_confirm, pattern=r"^chat:delete:\d+$"))
    app_.add_handler(CallbackQueryHandler(on_chat_delete_do,    pattern=r"^chat:delete:do:\d+$"))
    app_.add_handler(CallbackQueryHandler(on_chat_pin_toggle,   pattern=r"^chat:pin:\d+$"))
    app_.add_handler(CallbackQueryHandler(on_chat_share,        pattern=r"^chat:share:\d+$"))
    app_.add_handler(CallbackQueryHandler(on_chat_export_pdf,   pattern=r"^chat:export:pdf:\d+$"))
    app_.add_handler(CallbackQueryHandler(on_chat_export_md,    pattern=r"^chat:export:md:\d+$"))

    # error-handler
    app_.add_error_handler(on_error)

    return app_

@app.on_event("startup")
async def on_startup():
    global application, _public_url
    await init_db()

    application = build_application()
    await application.initialize()
    await application.start()

    _public_url = os.getenv("RENDER_EXTERNAL_URL") or os.getenv("PUBLIC_URL")
    if not _public_url:
        raise RuntimeError("Не найден PUBLIC_URL/RENDER_EXTERNAL_URL")

    webhook_url = f"{_public_url.rstrip('/')}/tg"
    await application.bot.set_webhook(webhook_url, max_connections=40, drop_pending_updates=False)
    logger.info("✅ Установлен Telegram webhook: %s", webhook_url)

    threading.Thread(target=_keepalive_loop, daemon=True).start()
    asyncio.get_event_loop().create_task(_webhook_guard_loop())
    asyncio.get_event_loop().create_task(_premium_expiry_notifier_loop())

    logger.info("🚀 Startup complete. Listening on port %s", PORT)

@app.on_event("shutdown")
async def on_shutdown():
    _keepalive_stop.set()
    try:
        if application is not None:
            # ВАЖНО: НЕ снимаем webhook — иначе Telegram перестанет будить Render!
            await application.stop()
            await application.shutdown()
    finally:
        logger.info("🛑 Shutdown complete")



# =========================
# Запуск
# =========================
if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=PORT)
